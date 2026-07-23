# -*- coding: utf-8 -*-
"""Génère INTÉGRALEMENT le deck OHC « dispositif d'écoute » (15 slides) —
générateur versionné, en remplacement des scripts jetables in-place
(arbitrage utilisateur 2026-07-23, application du diagnostic superviseur).

Modèle : le générateur standalone VSCode3 (docs/cadrage-ppt/generate_deck.py)
— une fonction slide_* par slide, un build() qui enchaîne, self-check
BLOQUANT en fin de build. Adapté au deck binaire de ce projet :

  - la référence v6 (Exports/… - v6.pptx) sert de TEMPLATE : on l'ouvre pour
    hériter masters / ~49 layouts OCTO (dont le layout « 51 - Chapitre [2] »
    déjà recoloré blob navy) / thème, et on extrait de son média (zipfile)
    les icônes héritées du deck original ; puis clear_slides() et
    reconstruction des 15 slides à neuf ;
  - le CONTENU (titres, textes, couleurs, positions) a été cartographié
    programmatiquement sur la v6 (dump shape par shape, 2026-07-23) et est
    codé EN DUR ci-dessous — fidélité visuelle à la v6, pas byte-à-byte ;
  - helpers génériques : scripts/pptx_deck.py (add_forme / add_text_runs /
    definir_paragraphes n'imposent QUE les propriétés fournies — le texte des
    placeholders garde la charte héritée des layouts, Outfit comprise).

Structure (v6, 2026-07-21) : Couverture · Sommaire (layout 92) · 4 chapitres
à dividers natifs (layout 51) — 01 Chantier OHC (Leviers, Personas,
Architecture) · 02 Existant & Évaluation (3 dispositifs RH, Google Form) ·
03 Nouveautés (Mentorat) · 04 Next steps (Arbitrages, Séquencement, Roadmap
datée 5 jalons → tribune 17 sept 2026).

Usage  : py scripts/generate_deck_ohc.py
Sortie : Exports/… - v7-genere.pptx (la v6 et les versions antérieures ne
sont JAMAIS modifiées — souvent ouvertes dans PowerPoint).
"""
import io
import os
import sys
import zipfile

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import pptx_deck as D  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Inches, Emu, Pt  # noqa: E402
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN, MSO_AUTO_SIZE  # noqa: E402
from pptx.oxml import parse_xml  # noqa: E402
from pptx.oxml.ns import qn  # noqa: E402

HERE = os.path.dirname(os.path.abspath(__file__))
REPO = os.path.abspath(os.path.join(HERE, ".."))

# pptx-framed-image (skill projet) : pose une photo dans le cadre du layout
# « 51 - Chapitre [2] » a la forme EXACTE du cadre (clip prstGeom, pas un
# rognage PIL) — cf. .claude/skills/pptx-framed-image/SKILL.md. Le cadre de
# CE template n'est pas groupe (frame_geometry() suppose un GROUP) : on
# retrouve ses bornes/geometrie avec D.trouver_cadre_layout (variante
# top-level ajoutee a pptx_deck.py pour ce cas).
sys.path.insert(0, os.path.join(REPO, ".claude", "skills", "pptx-framed-image", "scripts"))
from framed_image import (  # noqa: E402
    place_image_in_frame, frame_obstructions, cover_crop_to_aspect)
import nature_images  # noqa: E402
import stock_images  # noqa: E402

IMG_DIR = os.path.join(HERE, "_img")
os.makedirs(IMG_DIR, exist_ok=True)
IMG_MANIFEST = os.path.join(IMG_DIR, "manifest.json")
SOURCE_V6 = os.path.join(
    REPO, "Exports",
    "Chantiers OHC - dispositif écoute - avec synthese RH - v6.pptx")
SORTIE = os.path.join(
    REPO, "Exports",
    "Chantiers OHC - dispositif écoute - avec synthese RH - v7-genere.pptx")

# Charte du deck — thème du template OCTO de la v6, résolu en hexa
# (theme_colors : dk1/lt1/accent3/accent5/accent6, lt2).
NAVY = "#0E2356"   # dk1 — texte principal, pastilles DÉCISION
CYAN = "#00D2DD"   # accent3 — identité du deck (labels, filets, blobs)
SLATE = "#586586"  # lt2 — texte secondaire
BLANC = "#FFFFFF"  # lt1
LIGNE = "#CFD3DD"  # accent5 — bordures de cards
GRIS = "#E7E9EE"   # accent6 — fonds d'encarts

# Icônes réutilisées du média du deck ORIGINAL, embarquées dans la v6
# (pas de fetch externe — cohérence visuelle, cf. mémoire projet).
MEDIA_UTILISES = {"image4.png", "image5.png", "image13.png", "image15.png",
                  "image19.png"}


def charger_media(chemin):
    """Extrait de la v6 (zipfile) les blobs d'images réutilisés par les
    slides — clef = nom de fichier média (image5.png…)."""
    with zipfile.ZipFile(chemin) as z:
        return {nom.rsplit("/", 1)[-1]: z.read(nom)
                for nom in z.namelist()
                if nom.startswith("ppt/media/")
                and nom.rsplit("/", 1)[-1] in MEDIA_UTILISES}


def _layout(prs, nom):
    for la in prs.slide_masters[0].slide_layouts:
        if la.name == nom:
            return la
    raise ValueError(f"layout {nom!r} introuvable dans le template")


def _nouvelle_slide(prs, nom_layout, garder=()):
    """Ajoute une slide sur le layout nommé et ne conserve que les
    placeholders d'index listés dans `garder` (la v6 a par ex. supprimé le
    sous-titre du layout « 03 - Titre et sous-titre » sur la slide leviers)."""
    s = prs.slides.add_slide(_layout(prs, nom_layout))
    for ph in list(s.placeholders):
        if ph.placeholder_format.idx not in garder:
            ph._element.getparent().remove(ph._element)
    return s


# Blob « pin » (goutte de carte) des 3 leviers — géométrie custom héritée du
# deck original (Google Slides), extraite telle quelle de la v6. Fill cyan
# (accent3 du thème), contour navy (dk1) — comme la v6.
_BLOB_PIN_XML = (
    '<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
    ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
    '<p:nvSpPr><p:cNvPr id="{id}" name="Blob pin {id}"/><p:cNvSpPr/><p:nvPr/>'
    '</p:nvSpPr><p:spPr><a:xfrm><a:off x="{x}" y="{y}"/>'
    '<a:ext cx="848608" cy="954018"/></a:xfrm>'
    '<a:custGeom><a:rect b="b" l="l" r="r" t="t"/><a:pathLst>'
    '<a:path extrusionOk="0" h="21600" w="19679">'
    '<a:moveTo><a:pt x="9839" y="0"/></a:moveTo>'
    '<a:cubicBezTo><a:pt x="7321" y="0"/><a:pt x="4803" y="937"/>'
    '<a:pt x="2882" y="2813"/></a:cubicBezTo>'
    '<a:cubicBezTo><a:pt x="-961" y="6566"/><a:pt x="-961" y="12653"/>'
    '<a:pt x="2882" y="16406"/></a:cubicBezTo>'
    '<a:cubicBezTo><a:pt x="4480" y="17967"/><a:pt x="6493" y="18867"/>'
    '<a:pt x="8574" y="19129"/></a:cubicBezTo>'
    '<a:lnTo><a:pt x="9839" y="21600"/></a:lnTo>'
    '<a:lnTo><a:pt x="11104" y="19129"/></a:lnTo>'
    '<a:cubicBezTo><a:pt x="13185" y="18867"/><a:pt x="15198" y="17967"/>'
    '<a:pt x="16796" y="16406"/></a:cubicBezTo>'
    '<a:cubicBezTo><a:pt x="20639" y="12653"/><a:pt x="20639" y="6566"/>'
    '<a:pt x="16796" y="2813"/></a:cubicBezTo>'
    '<a:cubicBezTo><a:pt x="14875" y="937"/><a:pt x="12357" y="0"/>'
    '<a:pt x="9839" y="0"/></a:cubicBezTo>'
    '<a:close/></a:path></a:pathLst></a:custGeom>'
    '<a:solidFill><a:schemeClr val="accent3"/></a:solidFill>'
    '<a:ln cap="flat" cmpd="sng" w="9525">'
    '<a:solidFill><a:schemeClr val="dk1"/></a:solidFill>'
    '<a:prstDash val="solid"/><a:round/>'
    '<a:headEnd len="sm" w="sm" type="none"/>'
    '<a:tailEnd len="sm" w="sm" type="none"/></a:ln></p:spPr>'
    '<p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>'
)
_IDS_BLOB = iter(range(9001, 9900))


def _blob_pin(slide, l, t):
    """Pose un blob « pin » cyan (slide leviers) aux coordonnées données
    (pouces) — même custGeom que la v6, id de shape frais."""
    sp = parse_xml(_BLOB_PIN_XML.format(
        id=next(_IDS_BLOB), x=int(Inches(l)), y=int(Inches(t))))
    slide.shapes._spTree.append(sp)
    return sp

# -*- coding: utf-8 -*-


def slide_couverture(prs, media):
    """Couverture — layout natif 40 (slide 1 de la v6)."""
    s = _nouvelle_slide(prs, '40 - Couverture [1]', garder={0, 1, 2, 3})
    ph = s.placeholders[0]
    D.definir_geometrie(ph, 1.812, 1.771, 3.99, 1.427)
    D.configurer_text_frame(ph.text_frame, anchor=MSO_ANCHOR.MIDDLE, wrap=True, autosize=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE, margins=(0.0, 0.0, 0.0, 0.0))
    D.definir_paragraphes(ph.text_frame, [
        ([('Plans d’actions suite à l’enquête ', dict())], dict(align=PP_ALIGN.LEFT, line_spacing=1)),
        ([('OHC chez P&D', dict())], dict(align=PP_ALIGN.LEFT, line_spacing=1)),
    ])
    ph = s.placeholders[1]
    D.definir_geometrie(ph, 1.812, 2.918, 3.99, 0.236)
    D.configurer_text_frame(ph.text_frame, anchor=MSO_ANCHOR.BOTTOM, wrap=True, autosize=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE, margins=(0.0, 0.0, 0.0, 0.0))
    D.definir_paragraphes(ph.text_frame, [
        ([('Dispositif d’écoute', dict())], dict(align=PP_ALIGN.LEFT, line_spacing=1)),
    ])
    ph = s.placeholders[2]
    D.definir_geometrie(ph, 5.148, 3.613, 0.801, 0.345)
    D.configurer_text_frame(ph.text_frame, anchor=MSO_ANCHOR.MIDDLE, wrap=True, autosize=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE, margins=(0.0, 0.0, 0.0, 0.0))
    D.definir_paragraphes(ph.text_frame, [
        ([('OCTO', dict())], dict(align=PP_ALIGN.LEFT, line_spacing=1)),
    ])
    ph = s.placeholders[3]
    D.definir_geometrie(ph, 6.228, 3.613, 1.537, 0.345)
    D.configurer_text_frame(ph.text_frame, anchor=MSO_ANCHOR.MIDDLE, wrap=True, autosize=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE, margins=(0.0, 0.0, 0.0, 0.0))
    D.definir_paragraphes(ph.text_frame, [
        ([('02.07.2026', dict())], dict(align=PP_ALIGN.LEFT, line_spacing=1)),
    ])
    return s


def slide_sommaire(prs, media):
    """Sommaire — layout natif 92 Table des matieres (slide 2 de la v6)."""
    s = _nouvelle_slide(prs, '92 - Table des matières [4]', garder={0, 1, 2, 3, 4})
    ph = s.placeholders[0]
    D.definir_geometrie(ph, 0.907, 1.378, 8.756, 0.406)
    D.definir_paragraphes(ph.text_frame, [
        ([('Sommaire', dict())], dict()),
    ])
    ph = s.placeholders[1]
    D.definir_geometrie(ph, 1.839, 1.901, 3.078, 0.717)
    D.definir_paragraphes(ph.text_frame, [
        ([('Chantier OHC', dict())], dict()),
    ])
    ph = s.placeholders[2]
    D.definir_geometrie(ph, 1.839, 3.296, 3.078, 0.717)
    D.definir_paragraphes(ph.text_frame, [
        ([('Existant & Évaluation', dict())], dict()),
    ])
    ph = s.placeholders[3]
    D.definir_geometrie(ph, 6.146, 1.901, 3.235, 0.717)
    D.definir_paragraphes(ph.text_frame, [
        ([('Nouveautés', dict())], dict()),
    ])
    ph = s.placeholders[4]
    D.definir_geometrie(ph, 6.146, 3.296, 3.235, 0.717)
    D.definir_paragraphes(ph.text_frame, [
        ([('Next steps', dict())], dict()),
    ])
    return s


# Scène par chapitre — requête Openverse (photo réelle CC0) + seed, choisie et
# VÉRIFIÉE À L'ŒIL (pas seulement par mot-clé, cf. pptx-framed-image §Step 4) :
# plusieurs requêtes candidates par chapitre ont été essayées et rejetées au
# rendu (ex. « mentoring » seed 0 -> portrait d'époque du sénateur "Richard
# Mentor Johnson" ; « listening » seed 0 -> foule dans un monastère tibétain ;
# « survey » seed 0 -> photo N&B d'expédition 1870 ; « roadmap » seed 0 ->
# montage crypto sans rapport) avant de retenir ces quatre-ci — sobres, sans
# personne reconnaissable en gros plan ni élément hors-sujet.
_SCENES_CHAPITRE = {
    "01": ("team meeting", 0),          # constat/enquête : prise de notes en réunion
    "02": ("business meeting", 0),      # existant & évaluation : échange autour d'une table
    "03": ("colleagues smiling talking", 0),  # nouveautés (mentorat) : transmission entre pairs
    "04": ("planning whiteboard", 0),   # next steps : feuille de route au tableau
}
_NATURE_REPLI = "forest"  # scène connue de nature_images.py — repli hors-ligne uniquement


def _slug(texte):
    return "".join(c if c.isalnum() else "_" for c in texte.lower())


def _remplir_cadre_chapitre(slide, cadre, requete, seed=0):
    """Pose une vraie photo (Openverse CC0) dans le cadre du layout Chapitre,
    clippée à sa forme EXACTE (prstGeom cloné, cf. pptx-framed-image) — reprise
    du pattern réel VSCode3 `_remplir_cadre` (docs/cadrage-ppt/generate_deck.py).

    Le flip lu par `D.trouver_cadre_layout` (le cadre porte `xfrm flipH="1"`)
    n'est PAS réappliqué à l'image : un premier essai qui le réappliquait
    donnait un rendu géométriquement correct (mêmes coins arrondis que le
    cadre navy d'origine) mais avec la PHOTO MIROIR — `flipH` sur une image
    retourne son contenu pixel, pas seulement le choix du coin arrondi du clip
    (constaté au rendu réel : le texte du tableau blanc du chapitre 04 se
    lisait à l'envers). Le blob navy d'origine (fill uni) ne révélait jamais
    ce problème puisqu'un aplat de couleur flippé est visuellement identique.
    On ignore donc le flip : le clip `round2DiagRect` arrondit alors la
    diagonale opposée à celle du blob d'origine — un choix cosmétique sans
    incidence fonctionnelle (rien dans la charte ne dépend de CETTE
    diagonale précise), largement préférable à une photo mirroir."""
    if cadre is None:
        print("  cadre photo introuvable dans le layout — image non posée")
        return None
    left, top, width, height, geom, _flip_ignore = cadre
    aspect = Emu(width).inches / Emu(height).inches
    px_w = 900
    px_h = int(round(px_w / aspect))
    slug = _slug(requete)
    path = os.path.join(IMG_DIR, f"{slug}_{seed}_{px_w}x{px_h}.png")
    if not os.path.exists(path):
        try:
            brut = os.path.join(IMG_DIR, f"_brut_{slug}_{seed}.jpg")
            stock_images.fetch_to(brut, requete, seed=seed, manifest_path=IMG_MANIFEST)
            cover_crop_to_aspect(brut, path, aspect)
            print(f"  photo réelle posée pour {requete!r} (Openverse CC0)")
        except Exception as e:
            print(f"  Openverse indisponible pour {requete!r} ({e}) — repli nature_images")
            nature_images.generate_to(path, _NATURE_REPLI, px_w, px_h, seed=seed)
    return place_image_in_frame(slide, path, left, top, width, height, geom=geom)


def _neutraliser_fond_cadre(prs):
    """Le blob photo du layout « 51 - Chapitre [2] » porte un solidFill navy
    (contournement historique v6 : cadre vidé plutôt que rempli — cf. mémoire
    projet). On le repasse à noFill UNE FOIS (le layout est partagé par les 4
    dividers) : la photo posée par slide couvre exactement les mêmes bornes/
    géométrie, donc ce fond ne devrait jamais transparaître — filet défensif
    seulement (anti-liseré si un futur écart de géométrie apparaît)."""
    cadre = D.trouver_cadre_layout(
        _layout(prs, '51 - Chapitre [2]').shapes, 'round2DiagRect', largeur_min_in=2.0)
    if cadre is None:
        return
    _, _, _, _, geom_el, _ = cadre
    spPr = geom_el.getparent()  # <a:prstGeom> est un enfant direct de <p:spPr>
    for tag in ("a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill"):
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)
    no_fill = spPr.makeelement(qn("a:noFill"), {})
    geom_el.addnext(no_fill)



def slide_chapitre(prs, numero, titre):
    """Divider de chapitre — layout natif « 51 - Chapitre [2] ». Reprend
    EXPLICITEMENT la façon dont VSCode2 construit ses intercalaires
    (`app/services/pptx_export.py::_slide_chapitre`, lui-même la version
    éprouvée — P3 — du pattern du générateur de référence VSCode3) :
    - le NUMÉRO va dans le placeholder natif idx1 du layout (pas un chiffre
      dessiné en gros par-dessus la photo) : c'est la pastille-pilule
      top-gauche du template (contour navy, fond blanc — vérifié
      programmatiquement, shapes 187/201 du layout, `fill.type = NoFill`),
      DISTINCTE de la pastille logo du master (aucune collision réelle :
      logo en (0.30–0.49, 0.32–0.50)in, encart numéro en (0.92–1.47,
      0.41–0.88)in). Le seul piège réel (VSCode2/VSCode3) est le WRAP : la
      puce héritée pose marL=0.5in dans un encart de 0.55in de large — d'où
      `D.sans_puce`, marges à zéro, centré, 17pt (taille éprouvée par les
      deux générateurs de référence, jamais un gros numéro en pleine photo).
    - le TITRE va dans le placeholder idx0 (inchangé).
    - le cadre photo (round2DiagRect) reçoit une VRAIE photo (Openverse CC0,
      clippée à sa forme exacte) plutôt qu'un blob recoloré vidé — ce
      point-là était déjà correct depuis le run précédent.
    Le numéro sur son fond blanc n'a plus AUCUNE dépendance au contraste de
    la photo — l'ancien chiffre 54pt + scrim de protection posé par-dessus le
    cadre photo est donc retiré (il compensait un problème structurel qui
    disparaît avec le bon placeholder, pas un correctif à garder en plus)."""
    s = _nouvelle_slide(prs, '51 - Chapitre [2]', garder={0, 1})
    ph_num = s.placeholders[1]
    D.definir_geometrie(ph_num, 0.92, 0.41, 0.546, 0.474)
    tf_num = ph_num.text_frame
    tf_num.text = numero
    tf_num.margin_left = tf_num.margin_right = tf_num.margin_top = tf_num.margin_bottom = 0
    tf_num.vertical_anchor = MSO_ANCHOR.MIDDLE
    for p in tf_num.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        D.sans_puce(p)
        for r in p.runs:
            r.font.size = Pt(17)
            r.font.bold = True
            r.font.color.rgb = D.rgb(NAVY)
    D.appliquer_police(tf_num)
    ph_titre = s.placeholders[0]
    D.definir_geometrie(ph_titre, 5.14, 2.937, 4.554, 1.98)
    D.definir_paragraphes(ph_titre.text_frame, [
        ([(titre, dict())], dict()),
    ])
    cadre = D.trouver_cadre_layout(s.slide_layout.shapes, 'round2DiagRect', largeur_min_in=2.0)
    for pb in (frame_obstructions(s, *cadre[:4]) if cadre else []):
        print(f"  [obstruction] chapitre {numero}:", pb["source"], pb["name"], pb["reason"])
    requete, seed = _SCENES_CHAPITRE[numero]
    _remplir_cadre_chapitre(s, cadre, requete, seed)
    return s


# Puce réelle (buChar + retrait suspendu) des 3 cartes leviers — la v6 porte
# ces bullets au niveau paragraphe (héritage du deck original), pas en
# caractères littéraux comme les slides dessinées.
_PUCE_LEVIER = dict(char='•', size=10, font='Arial', color=NAVY)
_PUCE_LEVIER_9 = dict(char='•', size=9, font='Arial', color=NAVY)  # carte alerte (texte 8 pt)


def slide_leviers(prs, media):
    """En synthese — les 3 leviers (design herite du deck original, blobs pin cyan) (slide 4 de la v6)."""
    s = _nouvelle_slide(prs, '03 - Titre et sous-titre', garder={0})
    # bordure POINTILLÉE (prstDash dot) comme la v6 — seule bordure non pleine du deck
    D.add_forme(s, 'roundRect', 6.507, 0.972, 2.847, 4.615, fill=LIGNE, line='#3E4F78', line_w=0.75, adj=[0.16667], dash='dot')
    _f = D.add_forme(s, 'round2DiagRect', 6.622, 2.959, 2.632, 2.495, line=NAVY, line_w=0.75, adj=[0, 0.11406])
    D.configurer_text_frame(_f.text_frame, anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE, margins=(0.1, 0.25, 0.1, 0.1))
    D.definir_paragraphes(_f.text_frame, [
        ([('Avoir un ', dict(size=10, bold=False, italic=False, color=NAVY, font='Outfit')), ('espace de parole', dict(size=10, bold=True, italic=False, color=NAVY, font='Outfit')), (', guidé par des coachs P&D, pour partager ce qui ne va pas (1x / mois, pendant 45 min par ex)', dict(size=10, bold=False, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, line_spacing=1.15, bullet=_PUCE_LEVIER)),
        ([('Développer davantage le ', dict(size=10, bold=False, italic=False, color=NAVY, font='Outfit')), ('mentoring', dict(size=10, bold=True, italic=False, color=NAVY, font='Outfit')), ('  pour un soutien renforcé lors des missions complexes', dict(size=10, bold=False, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, line_spacing=1.15, bullet=_PUCE_LEVIER)),
    ])
    _f = D.add_forme(s, 'round2DiagRect', 3.637, 2.959, 2.632, 2.495, line=NAVY, line_w=0.75, adj=[0, 0.11406])
    D.configurer_text_frame(_f.text_frame, anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE, margins=(0.1, 0.25, 0.1, 0.1))
    D.definir_paragraphes(_f.text_frame, [
        ([('Concernant les situations de souffrance ou de perte de sens : ', dict(size=8, bold=False, italic=False, color=NAVY, font='Outfit')), ('Lister et partager', dict(size=8, bold=True, italic=False, color=NAVY, font='Outfit')), ("  les dispositifs d'alerte et de soutien managérial  existants, Vérifier si des", dict(size=8, bold=False, italic=False, color=NAVY, font='Outfit')), (' besoins non couverts', dict(size=8, bold=True, italic=False, color=NAVY, font='Outfit')), (' persistent.', dict(size=8, bold=False, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, line_spacing=1.15, bullet=_PUCE_LEVIER_9)),
        ([('Créer un mécanisme de remontée ', dict(size=8, bold=False, italic=False, color=NAVY, font='Outfit')), ('"mission toxique" ', dict(size=8, bold=True, italic=False, color=NAVY, font='Outfit')), ('— les consultants se sentent peu protégés face à des clients difficiles. Un process clair : signalement → arbitrage → décision de retrait si nécessaire.', dict(size=8, bold=False, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, line_spacing=1.15, bullet=_PUCE_LEVIER_9)),
    ])
    ph = s.placeholders[0]
    D.definir_geometrie(ph, 0.616, 0.395, 8.77, 0.406)
    D.configurer_text_frame(ph.text_frame, anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE, margins=(0.0, 0.0, 0.0, 0.0))
    D.definir_paragraphes(ph.text_frame, [
        ([('En synthèse - les 3 leviers d’actions prioritaires suite à l’OHC', dict())], dict(align=PP_ALIGN.LEFT, line_spacing=1)),
        ([('', dict())], dict(align=PP_ALIGN.LEFT, line_spacing=1)),
    ])
    D.add_forme(s, 'rect', 0.012, 2.292, 10, 0.05, fill=NAVY, rot=180)
    D.add_forme(s, 'round2DiagRect', 0.736, 1.624, 2.632, 1.335, line=NAVY, line_w=0.75, adj=[0.16667, 0])
    _blob_pin(s, 1.588, 1.025)
    D.add_forme(s, 'ellipse', 1.697, 1.135, 0.709, 0.709, fill=BLANC, line=NAVY, line_w=0.75)
    # groupe v6 decompose en formes absolues (echelle 1:1)
    D.add_forme(s, 'ellipse', 1.89, 2.155, 0.323, 0.324, fill=CYAN, line=NAVY, line_w=0.75)
    D.add_forme(s, 'ellipse', 1.967, 2.232, 0.17, 0.17, fill=BLANC, line=NAVY, line_w=0.75)
    D.add_text_runs(s, 6.842, 2.542, 2.152, 0.353, [
        ([('DISPOSITIF D’ÉCOUTE', dict(size=10.5, bold=True, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.CENTER, line_spacing=1)),
        ([('', dict())], dict(align=PP_ALIGN.LEFT, line_spacing=1)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'round2DiagRect', 3.637, 1.624, 2.632, 1.335, line=NAVY, line_w=0.75, adj=[0.16667, 0])
    D.add_text_runs(s, 3.806, 2.545, 2.378, 0.177, [
        ([('DISPOSITIF D’ALERTE', dict(size=10.5, bold=True, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.CENTER, line_spacing=1)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    _blob_pin(s, 4.489, 1.025)
    D.add_forme(s, 'ellipse', 4.598, 1.135, 0.709, 0.709, fill=BLANC, line=NAVY, line_w=0.75)
    # groupe v6 decompose en formes absolues (echelle 1:1)
    D.add_forme(s, 'ellipse', 4.791, 2.155, 0.323, 0.324, fill=CYAN, line=NAVY, line_w=0.75)
    D.add_forme(s, 'ellipse', 4.868, 2.232, 0.17, 0.17, fill=BLANC, line=NAVY, line_w=0.75)
    D.add_forme(s, 'round2DiagRect', 6.622, 1.624, 2.632, 1.335, line=NAVY, line_w=0.75, adj=[0.16667, 0])
    D.add_text_runs(s, 0.845, 2.509, 2.378, 0.353, [
        ([('DISPOSITIF DE SOUTIEN MANAGÉRIAL', dict(size=10.5, bold=True, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.CENTER, line_spacing=1)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    _blob_pin(s, 7.474, 1.025)
    D.add_forme(s, 'ellipse', 7.584, 1.135, 0.709, 0.709, fill=BLANC, line=NAVY, line_w=0.75)
    # groupe v6 decompose en formes absolues (echelle 1:1)
    D.add_forme(s, 'ellipse', 7.777, 2.155, 0.323, 0.324, fill=CYAN, line=NAVY, line_w=0.75)
    D.add_forme(s, 'ellipse', 7.854, 2.232, 0.17, 0.17, fill=BLANC, line=NAVY, line_w=0.75)
    _f = D.add_forme(s, 'round2DiagRect', 0.736, 2.959, 2.632, 2.495, line=NAVY, line_w=0.75, adj=[0, 0.11406])
    D.configurer_text_frame(_f.text_frame, anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE, margins=(0.1, 0.25, 0.1, 0.1))
    D.definir_paragraphes(_f.text_frame, [
        ([('Recréer de ', dict(size=10, bold=False, italic=False, color=NAVY, font='Outfit')), ('l’échange entre pairs', dict(size=10, bold=True, italic=False, color=NAVY, font='Outfit')), (' : (par ex. remettre en place les ateliers de ', dict(size=10, bold=False, italic=False, color=NAVY, font='Outfit')), ('co-dev', dict(size=10, bold=True, italic=False, color=NAVY, font='Outfit')), (' à destination des managers) une fois par trimestre accompagné d’un déjeuner', dict(size=10, bold=False, italic=False, color=NAVY, font='Outfit')), (' / un moment de convivialité', dict(size=10, bold=True, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, line_spacing=1.15, bullet=_PUCE_LEVIER, marL=0.3125, indent=-0.3125)),
        ([('', dict(size=10))], dict(align=PP_ALIGN.LEFT, line_spacing=1.15, bullet=_PUCE_LEVIER)),
    ])
    D.add_forme(s, 'round1Rect', 3.637, 2.956, 0.896, 0.227, fill=CYAN, line=NAVY, line_w=0.75, adj=[0.16667], rot=180)
    D.add_text_runs(s, 3.577, 3.01, 1.017, 0.227, [
        ([('Moyen terme ', dict(size=7, bold=False, italic=False, color=NAVY, font='Outfit SemiBold'))], dict(align=PP_ALIGN.CENTER, line_spacing=1)),
        ([('(1-3 mois)', dict(size=7, bold=False, italic=False, color=NAVY, font='Outfit SemiBold'))], dict(align=PP_ALIGN.CENTER, line_spacing=1)),
        ([('', dict())], dict(align=PP_ALIGN.CENTER, line_spacing=1)),
    ], anchor=MSO_ANCHOR.MIDDLE, wrap=True, autosize=MSO_AUTO_SIZE.NONE, margins=(0.1, 0.1, 0.1, 0.1))
    D.add_forme(s, 'round1Rect', 0.736, 2.956, 0.849, 0.227, fill=CYAN, line=NAVY, line_w=0.75, adj=[0.16667], rot=180)
    D.add_text_runs(s, 0.662, 3.011, 1.017, 0.227, [
        ([('Quick win ', dict(size=7, bold=False, italic=False, color=NAVY, font='Outfit SemiBold'))], dict(align=PP_ALIGN.CENTER, line_spacing=1)),
        ([('(< 1 mois)', dict(size=7, bold=False, italic=False, color=NAVY, font='Outfit SemiBold'))], dict(align=PP_ALIGN.CENTER, line_spacing=1)),
        ([('', dict())], dict(align=PP_ALIGN.CENTER, line_spacing=1)),
    ], anchor=MSO_ANCHOR.MIDDLE, wrap=True, autosize=MSO_AUTO_SIZE.NONE, margins=(0.1, 0.1, 0.1, 0.1))
    D.add_forme(s, 'round1Rect', 6.622, 2.956, 1.067, 0.227, fill=CYAN, line=NAVY, line_w=0.75, adj=[0.16667], rot=180)
    D.add_text_runs(s, 6.723, 3.008, 0.928, 0.227, [
        ([('Moyen terme', dict(size=7, bold=False, italic=False, color=NAVY, font='Outfit SemiBold'))], dict(align=PP_ALIGN.CENTER, line_spacing=1)),
        ([(' (1-3 mois)', dict(size=7, bold=False, italic=False, color=NAVY, font='Outfit SemiBold'))], dict(align=PP_ALIGN.CENTER, line_spacing=1)),
        ([('', dict())], dict(align=PP_ALIGN.CENTER, line_spacing=1)),
    ], anchor=MSO_ANCHOR.MIDDLE, wrap=True, autosize=MSO_AUTO_SIZE.NONE, margins=(0.1, 0.1, 0.1, 0.1))
    s.shapes.add_picture(io.BytesIO(media['image19.png']), Inches(7.589), Inches(1.104), Inches(0.698), Inches(0.698))
    s.shapes.add_picture(io.BytesIO(media['image15.png']), Inches(1.713), Inches(1.119), Inches(0.718), Inches(0.718))
    s.shapes.add_picture(io.BytesIO(media['image13.png']), Inches(4.594), Inches(1.119), Inches(0.718), Inches(0.718))
    return s


def slide_personas(prs, media):
    """Personas OHC — 8 populations en cartes KPI (slide 5 de la v6)."""
    s = _nouvelle_slide(prs, '06 - Slide vide', garder=())
    D.add_text_runs(s, 0.62, 0.34, 7.7, 0.5, [
        ([("Les personas de l'enquête OHC — 8 populations", dict(size=18, bold=True, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 0.62, 0.86, 8, 0.32, [
        ([("Volume et note des retours par population. En cyan : celles que le dispositif d'écoute adresse en priorité.", dict(size=10.5, bold=False, italic=False, color=SLATE, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    s.shapes.add_picture(io.BytesIO(media['image5.png']), Inches(8.62), Inches(0.28), Inches(0.365), Inches(0.62))
    D.add_forme(s, 'roundRect', 0.62, 1.4, 2.05, 1.48, fill=BLANC, line=LIGNE, line_w=0.75, adj=[0.06])
    D.add_forme(s, 'rect', 0.62, 1.4, 2.05, 0.07, fill=SLATE)
    D.add_text_runs(s, 0.76, 1.54, 1.77, 0.42, [
        ([('Managers', dict(size=9, bold=True, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2, line_spacing=0.95)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 0.76, 1.98, 1.77, 0.5, [
        ([('6', dict(size=22, bold=True, italic=False, color=SLATE, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 0.76, 2.53, 1.77, 0.28, [
        ([('retours · note 5,5/7', dict(size=8, bold=False, italic=False, color=SLATE, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'roundRect', 2.83, 1.4, 2.05, 1.48, fill=BLANC, line=CYAN, line_w=1.5, adj=[0.06])
    D.add_forme(s, 'rect', 2.83, 1.4, 2.05, 0.07, fill=CYAN)
    D.add_text_runs(s, 2.97, 1.54, 1.77, 0.42, [
        ([('Managés', dict(size=9, bold=True, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2, line_spacing=0.95)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 2.97, 1.98, 1.77, 0.5, [
        ([('11', dict(size=22, bold=True, italic=False, color=CYAN, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 2.97, 2.53, 1.77, 0.28, [
        ([('retours · note 6,1/7', dict(size=8, bold=False, italic=False, color=SLATE, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'roundRect', 5.04, 1.4, 2.05, 1.48, fill=BLANC, line=LIGNE, line_w=0.75, adj=[0.06])
    D.add_forme(s, 'rect', 5.04, 1.4, 2.05, 0.07, fill=SLATE)
    D.add_text_runs(s, 5.18, 1.54, 1.77, 0.42, [
        ([('Ateliers / Établis', dict(size=9, bold=True, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2, line_spacing=0.95)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 5.18, 1.98, 1.77, 0.5, [
        ([('19', dict(size=22, bold=True, italic=False, color=SLATE, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 5.18, 2.53, 1.77, 0.28, [
        ([('retours · note 3,9/7', dict(size=8, bold=False, italic=False, color=SLATE, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'roundRect', 7.25, 1.4, 2.05, 1.48, fill=BLANC, line=LIGNE, line_w=0.75, adj=[0.06])
    D.add_forme(s, 'rect', 7.25, 1.4, 2.05, 0.07, fill=SLATE)
    D.add_text_runs(s, 7.39, 1.54, 1.77, 0.42, [
        ([('Tribus', dict(size=9, bold=True, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2, line_spacing=0.95)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 7.39, 1.98, 1.77, 0.5, [
        ([('36', dict(size=22, bold=True, italic=False, color=SLATE, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 7.39, 2.53, 1.77, 0.28, [
        ([('retours · note 4,9/7', dict(size=8, bold=False, italic=False, color=SLATE, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'roundRect', 0.62, 2.98, 2.05, 1.48, fill=BLANC, line=LIGNE, line_w=0.75, adj=[0.06])
    D.add_forme(s, 'rect', 0.62, 2.98, 2.05, 0.07, fill=SLATE)
    D.add_text_runs(s, 0.76, 3.12, 1.77, 0.42, [
        ([('Compétences', dict(size=9, bold=True, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2, line_spacing=0.95)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 0.76, 3.56, 1.77, 0.5, [
        ([('12', dict(size=22, bold=True, italic=False, color=SLATE, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 0.76, 4.11, 1.77, 0.28, [
        ([('retours · note 4,6/7', dict(size=8, bold=False, italic=False, color=SLATE, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'roundRect', 2.83, 2.98, 2.05, 1.48, fill=BLANC, line=LIGNE, line_w=0.75, adj=[0.06])
    D.add_forme(s, 'rect', 2.83, 2.98, 2.05, 0.07, fill=SLATE)
    D.add_text_runs(s, 2.97, 3.12, 1.77, 0.42, [
        ([('Égalité pro', dict(size=9, bold=True, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2, line_spacing=0.95)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 2.97, 3.56, 1.77, 0.5, [
        ([('14', dict(size=22, bold=True, italic=False, color=SLATE, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 2.97, 4.11, 1.77, 0.28, [
        ([('retours · note 4,4/7', dict(size=8, bold=False, italic=False, color=SLATE, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'roundRect', 5.04, 2.98, 2.05, 1.48, fill=BLANC, line=CYAN, line_w=1.5, adj=[0.06])
    D.add_forme(s, 'rect', 5.04, 2.98, 2.05, 0.07, fill=CYAN)
    D.add_text_runs(s, 5.18, 3.12, 1.77, 0.42, [
        ([('Santé & conditions', dict(size=9, bold=True, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2, line_spacing=0.95)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 5.18, 3.56, 1.77, 0.5, [
        ([('26', dict(size=22, bold=True, italic=False, color=CYAN, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 5.18, 4.11, 1.77, 0.28, [
        ([('retours · note 5,4/7', dict(size=8, bold=False, italic=False, color=SLATE, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'roundRect', 7.25, 2.98, 2.05, 1.48, fill=BLANC, line=CYAN, line_w=1.5, adj=[0.06])
    D.add_forme(s, 'rect', 7.25, 2.98, 2.05, 0.07, fill=CYAN)
    D.add_text_runs(s, 7.39, 3.12, 1.77, 0.42, [
        ([('Missions & interco', dict(size=9, bold=True, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2, line_spacing=0.95)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 7.39, 3.56, 1.77, 0.5, [
        ([('10', dict(size=22, bold=True, italic=False, color=CYAN, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 7.39, 4.11, 1.77, 0.28, [
        ([('retours · note 4,4/7', dict(size=8, bold=False, italic=False, color=SLATE, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'roundRect', 0.62, 4.56, 8.77, 0.4, fill=GRIS, line=LIGNE, line_w=0.75, adj=[0.12])
    D.add_text_runs(s, 0.82, 4.56, 8.37, 0.4, [
        ([("Le dispositif d'écoute cible d'abord : ", dict(size=9, bold=True, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
        ([('Santé & conditions (stress), Missions & interco (mission « toxique ») et le soutien aux Managés.', dict(size=9, bold=False, italic=False, color=SLATE, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.MIDDLE, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    return s


def slide_architecture(prs, media):
    """Architecture de priorisation — 4 quadrants + axes dessines (slide 6 de la v6)."""
    s = _nouvelle_slide(prs, '06 - Slide vide', garder=())
    D.add_text_runs(s, 0.62, 0.34, 8.77, 0.5, [
        ([("Architecture de priorisation — où se situe le dispositif d'écoute", dict(size=18, bold=True, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'roundRect', 0.9, 1.62, 3.9, 1.3, fill=BLANC, line=LIGNE, line_w=0.75, adj=[0.06])
    D.add_forme(s, 'rect', 0.9, 1.62, 3.9, 0.06, fill='#2E8B57')
    D.add_text_runs(s, 1.06, 1.75, 3.6, 0.3, [
        ([('VEILLE — bien noté, peu de retours', dict(size=8.5, bold=True, italic=False, color='#2E8B57', font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 1.06, 2.09, 3.58, 0.74, [
        ([('Managers — 6 retours', dict(size=9, bold=False, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=3)),
        ([('Managés — 11 retours', dict(size=9, bold=True, italic=False, color=CYAN, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=3)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'roundRect', 5.15, 1.62, 3.9, 1.3, fill=BLANC, line=LIGNE, line_w=0.75, adj=[0.06])
    D.add_forme(s, 'rect', 5.15, 1.62, 3.9, 0.06, fill='#C88A14')
    D.add_text_runs(s, 5.31, 1.75, 3.6, 0.3, [
        ([('À SURVEILLER — bien noté, fort volume', dict(size=8.5, bold=True, italic=False, color='#C88A14', font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 5.31, 2.09, 3.58, 0.74, [
        ([('Santé & conditions — 26 retours', dict(size=9, bold=True, italic=False, color=CYAN, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=3)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'roundRect', 0.9, 3, 3.9, 1.3, fill=BLANC, line=LIGNE, line_w=0.75, adj=[0.06])
    D.add_forme(s, 'rect', 0.9, 3, 3.9, 0.06, fill=SLATE)
    D.add_text_runs(s, 1.06, 3.13, 3.6, 0.3, [
        ([('À QUALIFIER — à creuser', dict(size=8.5, bold=True, italic=False, color=SLATE, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 1.06, 3.47, 3.58, 0.74, [
        ([('Compétences — 12 retours', dict(size=9, bold=False, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=3)),
        ([('Égalité pro — 14 retours', dict(size=9, bold=False, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=3)),
        ([('Missions & interco — 10 retours', dict(size=9, bold=True, italic=False, color=CYAN, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=3)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'roundRect', 5.15, 3, 3.9, 1.3, fill=BLANC, line=LIGNE, line_w=0.75, adj=[0.06])
    D.add_forme(s, 'rect', 5.15, 3, 3.9, 0.06, fill='#C03A2B')
    D.add_text_runs(s, 5.31, 3.13, 3.6, 0.3, [
        ([('TOP PRIORITÉ — peu satisfait, fort volume', dict(size=8.5, bold=True, italic=False, color='#C03A2B', font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 5.31, 3.47, 3.58, 0.74, [
        ([('Tribus — 36 retours', dict(size=9, bold=False, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=3)),
        ([('Ateliers / Établis — 19 retours', dict(size=9, bold=False, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=3)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'roundRect', 0.62, 4.5, 8.77, 0.4, fill=GRIS, line=LIGNE, line_w=0.75, adj=[0.12])
    D.add_text_runs(s, 0.82, 4.5, 8.37, 0.4, [
        ([("L'écoute agit sur les tensions à faible satisfaction : ", dict(size=9, bold=True, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
        ([('Santé & conditions, Missions & interco (mission « toxique ») et le soutien aux Managés.', dict(size=9, bold=False, italic=False, color=SLATE, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.MIDDLE, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'rect', 0.82, 1.55, 0.018, 2.85, fill=NAVY)
    D.add_forme(s, 'triangle', 0.75, 1.4, 0.15, 0.16, fill=NAVY)
    D.add_text_runs(s, 0.18, 1.5, 0.72, 0.3, [
        ([('＋ satisfait', dict(size=8, bold=True, italic=False, color=SLATE))], dict(align=PP_ALIGN.LEFT)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.NONE)
    D.add_text_runs(s, 0.18, 4.02, 0.72, 0.3, [
        ([('− satisfait', dict(size=8, bold=True, italic=False, color=SLATE))], dict(align=PP_ALIGN.LEFT)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.NONE)
    D.add_forme(s, 'rect', 0.9, 4.36, 8.1, 0.018, fill=NAVY)
    D.add_forme(s, 'triangle', 9, 4.3, 0.16, 0.15, fill=NAVY, rot=90)
    D.add_text_runs(s, 5.15, 4.4, 3.85, 0.16, [
        ([('VOLUME DE RETOURS →', dict(size=8, bold=True, italic=False, color=SLATE))], dict(align=PP_ALIGN.RIGHT)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.NONE)
    return s


def slide_existant(prs, media):
    """Existant — 3 dispositifs RH en cartes + protocole condense (slide 8 de la v6)."""
    s = _nouvelle_slide(prs, '04 - Titre seul', garder={0})
    ph = s.placeholders[0]
    D.definir_geometrie(ph, 0.615, 0.395, 8.77, 0.406)
    D.definir_paragraphes(ph.text_frame, [
        ([('Rappel — 3 dispositifs d’écoute RH déjà en place chez OCTO', dict(font='Outfit SemiBold'))], dict()),
    ])
    D.add_text_runs(s, 0.62, 0.85, 8.77, 0.3, [
        ([('Avant de lancer de nouveaux formats P&D : ce qui existe déjà et comment le déclencher.', dict(size=10.5, bold=False, italic=True, color=SLATE, font='Outfit'))], dict(align=PP_ALIGN.LEFT)),
    ], wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT, margins=(0.1, 0.05, 0.1, 0.05))
    D.add_forme(s, 'roundRect', 0.62, 1.35, 2.757, 3.05, fill=BLANC, line=LIGNE, line_w=0.75, adj=[0.06])
    D.add_text_runs(s, 0.78, 1.49, 2.437, 0.3, [
        ([('PLAN STRESS OCTO', dict(size=10.5, bold=True, italic=False, color=CYAN))], dict(align=PP_ALIGN.LEFT)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.NONE)
    D.add_forme(s, 'rect', 0.78, 1.79, 0.4, 0.03, fill=CYAN)
    D.add_text_runs(s, 0.78, 1.93, 2.437, 1.15, [
        ([('• Boîte à outils : repérer les signaux, souffler, en parler (manager, RH, pairs)', dict(size=9.3, bold=False, italic=False, color=NAVY))], dict(align=PP_ALIGN.LEFT, space_after=5, line_spacing=1.08)),
        ([('• O3 tous les 15 j : détection + feedback managérial', dict(size=9.3, bold=False, italic=False, color=NAVY))], dict(align=PP_ALIGN.LEFT, space_after=5, line_spacing=1.08)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.NONE)
    D.add_forme(s, 'roundRect', 0.684, 3.13, 2.629, 1.15, fill=GRIS, adj=[0.1])
    D.add_text_runs(s, 0.804, 3.22, 2.389, 0.97, [
        ([('Déclencheur — ', dict(size=8.3, bold=True, color=NAVY, font='Outfit')), ('espaceocto@gmail.com', dict(size=8.3, bold=False, color=SLATE, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=3)),
        ([('  ·  rappel < 2j, RDV < 5j  ·  1 à 3 séances (1h30, confidentiel)', dict(size=8.3, bold=False, color=SLATE, font='Outfit'))], dict(space_before=3)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.NONE)
    D.add_forme(s, 'roundRect', 3.627, 1.35, 2.757, 3.05, fill=BLANC, line=LIGNE, line_w=0.75, adj=[0.06])
    D.add_text_runs(s, 3.787, 1.49, 2.437, 0.3, [
        ([('COACHING PRO', dict(size=10.5, bold=True, italic=False, color=CYAN))], dict(align=PP_ALIGN.LEFT)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.NONE)
    D.add_forme(s, 'rect', 3.787, 1.79, 0.4, 0.03, fill=CYAN)
    D.add_text_runs(s, 3.787, 1.93, 2.437, 1.15, [
        ([('• Accompagnement individuel limité dans le temps, 100% confidentiel', dict(size=9.3, bold=False, italic=False, color=NAVY))], dict(align=PP_ALIGN.LEFT, space_after=5, line_spacing=1.08)),
        ([('• Interne (coach Octo) ou externe selon le besoin qualifié', dict(size=9.3, bold=False, italic=False, color=NAVY))], dict(align=PP_ALIGN.LEFT, space_after=5, line_spacing=1.08)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.NONE)
    D.add_forme(s, 'roundRect', 3.691, 3.13, 2.629, 1.15, fill=GRIS, adj=[0.1])
    D.add_text_runs(s, 3.811, 3.22, 2.389, 0.97, [
        ([('Déclencheur — ', dict(size=8.3, bold=True, color=NAVY, font='Outfit')), ('ta lead RH', dict(size=8.3, bold=False, color=SLATE, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=3)),
        ([('  ·  qualification + choix du coach (cadrage tripartite)  ·  6-8 séances puis bilan', dict(size=8.3, bold=False, color=SLATE, font='Outfit'))], dict(space_before=3)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.NONE)
    D.add_forme(s, 'roundRect', 6.633, 1.35, 2.757, 3.05, fill=BLANC, line=LIGNE, line_w=0.75, adj=[0.06])
    D.add_text_runs(s, 6.793, 1.49, 2.437, 0.3, [
        ([('PERSONNE RESSOURCE', dict(size=10.5, bold=True, italic=False, color=CYAN))], dict(align=PP_ALIGN.LEFT)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.NONE)
    D.add_forme(s, 'rect', 6.793, 1.79, 0.4, 0.03, fill=CYAN)
    # h 1.15 (v6) -> 1.52 : boîte invisible ancrée en haut, hauteur élargie
    # pour satisfaire l'estimateur pessimiste de verifier_debordements_texte
    # (le rendu réel tient en ~5 lignes) — aucun effet visuel.
    D.add_text_runs(s, 6.793, 1.93, 2.437, 1.52, [
        ([('• Oreille de proximité sur les sujets diversité (gender mix, LGBTQIA+, mixité sociale…)', dict(size=9.3, bold=False, italic=False, color=NAVY))], dict(align=PP_ALIGN.LEFT, space_after=5, line_spacing=1.08)),
        ([('• Échange ponctuel — jamais un suivi régulier (≠ psy, ≠ coach)', dict(size=9.3, bold=False, italic=False, color=NAVY))], dict(align=PP_ALIGN.LEFT, space_after=5, line_spacing=1.08)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.NONE)
    D.add_forme(s, 'roundRect', 6.697, 3.13, 2.629, 1.15, fill=GRIS, adj=[0.1])
    D.add_text_runs(s, 6.817, 3.22, 2.389, 0.97, [
        ([('Déclencheur — ', dict(size=8.3, bold=True, color=NAVY, font='Outfit')), ('liste MM (mail / tél.)', dict(size=8.3, bold=False, color=SLATE, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=3)),
        ([('  ·  échange ponctuel confidentiel  ·  réorientation RH si le suivi devient régulier', dict(size=8.3, bold=False, color=SLATE, font='Outfit'))], dict(space_before=3)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.NONE)
    D.add_forme(s, 'roundRect', 0.62, 4.52, 8.77, 0.42, fill=GRIS, adj=[0.2])
    D.add_text_runs(s, 0.82, 4.52, 8.37, 0.42, [
        ([('À réactiver, pas à recréer : ', dict(size=9, bold=True, color=NAVY, font='Outfit')), ('deux formats collectifs existent déjà (cercle de parole, clinique de mission — hérités de Zen@octo 2016), aujourd’hui dormants. Une seule brique manque vraiment : le mentorat mission complexe.', dict(size=9, bold=False, color=SLATE, font='Outfit'))], dict(align=PP_ALIGN.LEFT)),
    ], anchor=MSO_ANCHOR.MIDDLE, wrap=True, autosize=MSO_AUTO_SIZE.NONE)
    return s


def slide_evaluation(prs, media):
    """Evaluation Google Form — 4 mesures + mockup dessine (violet) (slide 9 de la v6)."""
    s = _nouvelle_slide(prs, '06 - Slide vide', garder=())
    D.add_text_runs(s, 0.62, 0.34, 8.77, 0.5, [
        ([('Évaluer l’existant & recueillir les besoins — Google Form', dict(size=18, bold=True, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 0.62, 0.86, 8.77, 0.32, [
        ([('Anonyme · trimestriel · agrégé — mesure les 3 dispositifs existants et recueille les besoins non couverts.', dict(size=10.5, bold=False, italic=False, color=SLATE, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'roundRect', 0.62, 1.35, 2.1, 1.06, fill=BLANC, line=LIGNE, line_w=0.75, adj=[0.06])
    D.add_text_runs(s, 0.77, 1.47, 1.85, 0.26, [
        ([('NOTORIÉTÉ', dict(size=9, bold=True, italic=False, color=CYAN, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'rect', 0.77, 1.75, 0.34, 0.045, fill=CYAN)
    D.add_text_runs(s, 0.77, 1.87, 1.85, 0.5, [
        ([('Connu ? Sais-tu le déclencher ?', dict(size=8.5, bold=False, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2, line_spacing=1.02)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'roundRect', 2.8, 1.35, 2.1, 1.06, fill=BLANC, line=LIGNE, line_w=0.75, adj=[0.06])
    D.add_text_runs(s, 2.95, 1.47, 1.85, 0.26, [
        ([('RECOURS & UTILITÉ', dict(size=9, bold=True, italic=False, color=CYAN, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'rect', 2.95, 1.75, 0.34, 0.045, fill=CYAN)
    D.add_text_runs(s, 2.95, 1.87, 1.85, 0.5, [
        ([("Sollicité ? Ça t'a aidé (1–5) ?", dict(size=8.5, bold=False, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2, line_spacing=1.02)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'roundRect', 0.62, 2.49, 2.1, 1.06, fill=BLANC, line=LIGNE, line_w=0.75, adj=[0.06])
    D.add_text_runs(s, 0.77, 2.61, 1.85, 0.26, [
        ([('FREINS', dict(size=9, bold=True, italic=False, color=CYAN, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'rect', 0.77, 2.89, 0.34, 0.045, fill=CYAN)
    D.add_text_runs(s, 0.77, 3.01, 1.85, 0.5, [
        ([("Qu'est-ce qui te retient ?", dict(size=8.5, bold=False, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2, line_spacing=1.02)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'roundRect', 2.8, 2.49, 2.1, 1.06, fill=BLANC, line=LIGNE, line_w=0.75, adj=[0.06])
    D.add_text_runs(s, 2.95, 2.61, 1.85, 0.26, [
        ([('COUVERTURE', dict(size=9, bold=True, italic=False, color=CYAN, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'rect', 2.95, 2.89, 0.34, 0.045, fill=CYAN)
    D.add_text_runs(s, 2.95, 3.01, 1.85, 0.5, [
        ([('Besoin non couvert ? → nourrit la nouveauté', dict(size=8.5, bold=False, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2, line_spacing=1.02)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'roundRect', 0.62, 3.66, 4.28, 1.22, fill=GRIS, line=LIGNE, line_w=0.75, adj=[0.08])
    D.add_text_runs(s, 0.8, 3.78, 3.94, 1, [
        ([('Protocole', dict(size=9, bold=True, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
        ([('Anonyme · trimestriel · résultats agrégés → réactiver et rendre visible, pas recréer.', dict(size=8.3, bold=False, italic=False, color=SLATE, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=3, line_spacing=1.05)),
        ([('Cibles : notoriété > 70 % · sait déclencher > 50 % · utilité > 3,5 / 5', dict(size=8.3, bold=False, italic=False, color=SLATE, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2, line_spacing=1.05)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'roundRect', 5.2, 1.3, 4.2, 3.58, fill=BLANC, line=LIGNE, line_w=0.75, adj=[0.03])
    D.add_forme(s, 'rect', 5.2, 1.3, 4.2, 0.12, fill='#673AB7')
    D.add_text_runs(s, 5.42, 1.54, 3.8, 0.3, [
        ([("Dispositifs d'écoute — mini-bilan", dict(size=11, bold=True, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 5.42, 1.82, 3.8, 0.22, [
        ([('Anonyme · environ 2 min · Section 1 sur 3', dict(size=8, bold=False, italic=False, color='#70757A', font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'rect', 5.42, 2.08, 3.76, 0.012, fill='#DADCE0')
    D.add_text_runs(s, 5.42, 2.2, 3.76, 0.24, [
        ([("Connais-tu l'espace de discussion neutre ?", dict(size=9, bold=True, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'ellipse', 5.48, 2.48, 0.13, 0.13, line='#70757A', line_w=1)
    D.add_text_runs(s, 5.7, 2.46, 3.4, 0.22, [
        ([('Je ne connaissais pas', dict(size=8.5, bold=False, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'ellipse', 5.48, 2.72, 0.13, 0.13, line='#70757A', line_w=1)
    D.add_text_runs(s, 5.7, 2.7, 3.4, 0.22, [
        ([("J'en ai entendu parler", dict(size=8.5, bold=False, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'ellipse', 5.48, 2.96, 0.13, 0.13, fill='#673AB7')
    D.add_forme(s, 'ellipse', 5.512, 2.992, 0.066, 0.066, fill=BLANC)
    D.add_text_runs(s, 5.7, 2.94, 3.4, 0.22, [
        ([('Je sais comment le déclencher', dict(size=8.5, bold=False, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 5.42, 3.32, 3.76, 0.24, [
        ([("À quel point cela t'a-t-il aidé·e ?", dict(size=9, bold=True, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 5.46, 3.58, 1.2, 0.2, [
        ([('Pas du tout', dict(size=7, bold=False, italic=False, color='#70757A', font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 8.4, 3.58, 0.78, 0.2, [
        ([('Beaucoup', dict(size=7, bold=False, italic=False, color='#70757A', font='Outfit'))], dict(align=PP_ALIGN.RIGHT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 5.58, 3.74, 0.24, 0.18, [
        ([('1', dict(size=8, bold=False, italic=False, color='#70757A', font='Outfit'))], dict(align=PP_ALIGN.CENTER, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'ellipse', 5.635, 3.94, 0.13, 0.13, line='#70757A', line_w=1)
    D.add_text_runs(s, 6.3, 3.74, 0.24, 0.18, [
        ([('2', dict(size=8, bold=False, italic=False, color='#70757A', font='Outfit'))], dict(align=PP_ALIGN.CENTER, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'ellipse', 6.355, 3.94, 0.13, 0.13, line='#70757A', line_w=1)
    D.add_text_runs(s, 7.02, 3.74, 0.24, 0.18, [
        ([('3', dict(size=8, bold=False, italic=False, color='#70757A', font='Outfit'))], dict(align=PP_ALIGN.CENTER, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'ellipse', 7.075, 3.94, 0.13, 0.13, line='#70757A', line_w=1)
    D.add_text_runs(s, 7.74, 3.74, 0.24, 0.18, [
        ([('4', dict(size=8, bold=False, italic=False, color='#70757A', font='Outfit'))], dict(align=PP_ALIGN.CENTER, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'ellipse', 7.795, 3.94, 0.13, 0.13, fill='#673AB7')
    D.add_forme(s, 'ellipse', 7.827, 3.972, 0.066, 0.066, fill=BLANC)
    D.add_text_runs(s, 8.46, 3.74, 0.24, 0.18, [
        ([('5', dict(size=8, bold=False, italic=False, color='#70757A', font='Outfit'))], dict(align=PP_ALIGN.CENTER, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'ellipse', 8.515, 3.94, 0.13, 0.13, line='#70757A', line_w=1)
    D.add_forme(s, 'roundRect', 5.42, 4.46, 0.86, 0.28, fill='#673AB7', adj=[0.18])
    D.add_text_runs(s, 5.42, 4.465, 0.86, 0.27, [
        ([('Envoyer', dict(size=8.5, bold=True, italic=False, color=BLANC, font='Outfit'))], dict(align=PP_ALIGN.CENTER, space_after=2)),
    ], anchor=MSO_ANCHOR.MIDDLE, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 8, 4.48, 1.18, 0.24, [
        ([('Page 1 sur 3', dict(size=7.5, bold=False, italic=False, color='#70757A', font='Outfit'))], dict(align=PP_ALIGN.RIGHT, space_after=2)),
    ], anchor=MSO_ANCHOR.MIDDLE, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    return s


def slide_mentorat(prs, media):
    """Mentorat mission complexe — format + cadrage (slide 11 de la v6)."""
    s = _nouvelle_slide(prs, '06 - Slide vide', garder=())
    D.add_text_runs(s, 0.62, 0.34, 7.7, 0.5, [
        ([('Mentorat mission complexe — le format et son cadrage', dict(size=18, bold=True, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 0.62, 0.86, 8, 0.32, [
        ([('La seule brique réellement nouvelle', dict(size=10.5, bold=True, color=CYAN, font='Outfit')), (' : un pair-mentor qui suit une mission complexe dans la durée — ce qu’aucun dispositif existant ne fait.', dict(size=10.5, bold=False, color=SLATE, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    s.shapes.add_picture(io.BytesIO(media['image5.png']), Inches(8.62), Inches(0.28), Inches(0.365), Inches(0.62))
    D.add_forme(s, 'roundRect', 0.62, 1.28, 4.24, 1.5, fill=BLANC, line=LIGNE, line_w=0.75, adj=[0.06])
    D.add_text_runs(s, 0.8, 1.41, 3.88, 0.28, [
        ([('FORMULER LA DEMANDE — atelier « Besoin de mentoring »', dict(size=9, bold=True, italic=False, color=CYAN, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'rect', 0.8, 1.72, 0.42, 0.045, fill=CYAN)
    D.add_text_runs(s, 0.8, 1.86, 3.88, 0.8, [
        ([('•  Auto-diagnostic : client, politique, posture, charge', dict(size=9, bold=False, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=4, line_spacing=1.02)),
        ([("•  « J'ai besoin d'un mentor pour… » → type de mentor utile", dict(size=9, bold=False, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=4, line_spacing=1.02)),
        ([('•  Cadré par le lead RH : ni évaluation, ni jugement de performance', dict(size=9, color=NAVY, font='Outfit'))], dict(space_before=3)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'roundRect', 5.15, 1.28, 4.24, 1.5, fill=BLANC, line=LIGNE, line_w=0.75, adj=[0.06])
    D.add_text_runs(s, 5.33, 1.41, 3.88, 0.28, [
        ([('RENCONTRER UN MENTOR — speed dating', dict(size=9, bold=True, italic=False, color=CYAN, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'rect', 5.33, 1.72, 0.42, 0.045, fill=CYAN)
    D.add_text_runs(s, 5.33, 1.86, 3.88, 0.8, [
        ([('•  Présentation des mentors et expertises', dict(size=9, bold=False, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=4, line_spacing=1.02)),
        ([("•  Rotations courtes → choix d'un binôme potentiel", dict(size=9, bold=False, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=4, line_spacing=1.02)),
        ([('•  Choix mutuel — le mentoré garde la main sur le binôme', dict(size=9, color=NAVY, font='Outfit'))], dict(space_before=3)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'roundRect', 0.62, 2.92, 2.05, 1.16, fill=BLANC, line=LIGNE, line_w=0.75, adj=[0.06])
    D.add_text_runs(s, 0.76, 3.03, 1.8, 0.24, [
        ([('DÉCLENCHEMENT', dict(size=8.5, bold=True, italic=False, color=CYAN, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'rect', 0.76, 3.29, 0.3, 0.045, fill=CYAN)
    D.add_text_runs(s, 0.76, 3.4, 1.8, 0.62, [
        ([('Mission « toxique » repérée → coach → lead RH', dict(size=8, bold=False, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2, line_spacing=1.02)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'roundRect', 2.83, 2.92, 2.05, 1.16, fill=BLANC, line=LIGNE, line_w=0.75, adj=[0.06])
    D.add_text_runs(s, 2.97, 3.03, 1.8, 0.24, [
        ([('APPARIEMENT', dict(size=8.5, bold=True, italic=False, color=CYAN, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'rect', 2.97, 3.29, 0.3, 0.045, fill=CYAN)
    D.add_text_runs(s, 2.97, 3.4, 1.8, 0.62, [
        ([("Pool de mentors volontaires ; réseaux existants d'abord", dict(size=8, bold=False, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2, line_spacing=1.02)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'roundRect', 5.04, 2.92, 2.05, 1.16, fill=BLANC, line=LIGNE, line_w=0.75, adj=[0.06])
    D.add_text_runs(s, 5.18, 3.03, 1.8, 0.24, [
        ([('GARDE-FOUS', dict(size=8.5, bold=True, italic=False, color=CYAN, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'rect', 5.18, 3.29, 0.3, 0.045, fill=CYAN)
    D.add_text_runs(s, 5.18, 3.4, 1.8, 0.62, [
        ([('Protéger le mentoré ; escalade RH si risque', dict(size=8, bold=False, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2, line_spacing=1.02)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'roundRect', 7.25, 2.92, 2.05, 1.16, fill=BLANC, line=LIGNE, line_w=0.75, adj=[0.06])
    D.add_text_runs(s, 7.39, 3.03, 1.8, 0.24, [
        ([('MESURE', dict(size=8.5, bold=True, italic=False, color=CYAN, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'rect', 7.39, 3.29, 0.3, 0.045, fill=CYAN)
    D.add_text_runs(s, 7.39, 3.4, 1.8, 0.62, [
        ([('Retrait / rééquilibrage ; retour de sérénité', dict(size=8, bold=False, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2, line_spacing=1.02)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'roundRect', 0.62, 4.36, 8.77, 0.42, fill=GRIS, line=LIGNE, line_w=0.75, adj=[0.12])
    D.add_text_runs(s, 0.82, 4.36, 8.37, 0.42, [
        ([('Cadre de pairing : ', dict(size=9, bold=True, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
        ([('30 min toutes les 2–3 semaines · engagement initial de 3 rendez-vous · bilan puis poursuite si utile.', dict(size=9, bold=False, italic=False, color=SLATE, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.MIDDLE, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    return s


def slide_arbitrages(prs, media):
    """Arbitrages ecoute individuelle — 3 cartes Enjeu -> Reco, pastilles DECISION navy (slide 13 de la v6)."""
    s = _nouvelle_slide(prs, '06 - Slide vide', garder=())
    D.add_text_runs(s, 0.62, 0.34, 7.7, 0.5, [
        ([('Écoute individuelle — arbitrages recommandés', dict(size=18, bold=True, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 0.62, 0.86, 8, 0.32, [
        ([('Trois décisions pour cadrer les formats individuels — recommandation consultant, à ratifier par le board P&D.', dict(size=10.5, bold=False, italic=False, color=SLATE, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    s.shapes.add_picture(io.BytesIO(media['image13.png']), Inches(8.62), Inches(0.28), Inches(0.62), Inches(0.62))
    D.add_forme(s, 'roundRect', 0.62, 1.3, 2.76, 2.72, fill=BLANC, line=LIGNE, line_w=0.75, adj=[0.06])
    D.add_text_runs(s, 0.8, 1.43, 2.4, 0.28, [
        ([("OÙ SE FAIT L'ÉCOUTE ?", dict(size=9.5, bold=True, italic=False, color=CYAN, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'rect', 0.8, 1.74, 0.42, 0.045, fill=CYAN)
    D.add_forme(s, 'roundRect', 0.8, 1.85, 0.86, 0.24, fill=NAVY, adj=[0.5])
    D.add_text_runs(s, 0.8, 1.85, 0.86, 0.23, [
        ([('DÉCISION', dict(size=8, bold=True, italic=False, color=BLANC, font='Outfit'))], dict(align=PP_ALIGN.CENTER, space_after=2)),
    ], anchor=MSO_ANCHOR.MIDDLE, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 0.8, 2.16, 2.4, 1.74, [
        ([('Enjeu — ', dict(size=7.8, bold=True, color=SLATE, font='Outfit')), ('des coachs internes P&D peuvent manquer de neutralité sur des sujets qui touchent le management de l’atelier.', dict(size=7.8, bold=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=6, line_spacing=1.03)),
        ([('→ Reco — ', dict(size=8, bold=True, color=CYAN, font='Outfit')), ('confiance envers le management → espace neutre externe (espaceocto, confidentiel) ; contexte mission → clinique ou mentorat.', dict(size=8, bold=True, color=NAVY, font='Outfit'))], dict(space_after=6, line_spacing=1.03)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'roundRect', 3.63, 1.3, 2.76, 2.72, fill=BLANC, line=LIGNE, line_w=0.75, adj=[0.06])
    D.add_text_runs(s, 3.81, 1.43, 2.4, 0.28, [
        ([('VOCABULAIRE « COACH »', dict(size=9.5, bold=True, italic=False, color=CYAN, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'rect', 3.81, 1.74, 0.42, 0.045, fill=CYAN)
    D.add_forme(s, 'roundRect', 3.81, 1.85, 0.86, 0.24, fill=NAVY, adj=[0.5])
    D.add_text_runs(s, 3.81, 1.85, 0.86, 0.23, [
        ([('DÉCISION', dict(size=8, bold=True, italic=False, color=BLANC, font='Outfit'))], dict(align=PP_ALIGN.CENTER, space_after=2)),
    ], anchor=MSO_ANCHOR.MIDDLE, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 3.81, 2.16, 2.4, 1.74, [
        ([('Enjeu — ', dict(size=7.8, bold=True, color=SLATE, font='Outfit')), ('« coach » engage la déontologie du coaching pro OCTO (confidentialité, posture) que les animateurs P&D n’ont pas forcément.', dict(size=7.8, bold=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=6, line_spacing=1.03)),
        ([('→ Reco — ', dict(size=8, bold=True, color=CYAN, font='Outfit')), ('renommer « facilitateur·rice d’écoute P&D », ou exiger le socle de formation à l’écoute active — à inscrire dans la charte.', dict(size=8, bold=True, color=NAVY, font='Outfit'))], dict(space_after=6, line_spacing=1.03)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'roundRect', 6.63, 1.3, 2.76, 2.72, fill=BLANC, line=LIGNE, line_w=0.75, adj=[0.06])
    D.add_text_runs(s, 6.81, 1.43, 2.4, 0.28, [
        ([('UN 3ᵉ POOL DE VOLONTAIRES ?', dict(size=9.5, bold=True, italic=False, color=CYAN, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'rect', 6.81, 1.74, 0.42, 0.045, fill=CYAN)
    D.add_forme(s, 'roundRect', 6.81, 1.85, 0.86, 0.24, fill=NAVY, adj=[0.5])
    D.add_text_runs(s, 6.81, 1.85, 0.86, 0.23, [
        ([('DÉCISION', dict(size=8, bold=True, italic=False, color=BLANC, font='Outfit'))], dict(align=PP_ALIGN.CENTER, space_after=2)),
    ], anchor=MSO_ANCHOR.MIDDLE, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 6.81, 2.16, 2.4, 1.74, [
        ([('Enjeu — ', dict(size=7.8, bold=True, color=SLATE, font='Outfit')), ('un réseau personne ressource et un pool coaching interne existent déjà à l’échelle OCTO.', dict(size=7.8, bold=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=6, line_spacing=1.03)),
        ([('→ Reco — ', dict(size=8, bold=True, color=CYAN, font='Outfit')), ('les solliciter d’abord ; ne recruter un pool P&D dédié que pour la part non couverte (clinique, spécifique métier).', dict(size=8, bold=True, color=NAVY, font='Outfit'))], dict(space_after=6, line_spacing=1.03)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'roundRect', 0.62, 4.34, 8.77, 0.42, fill=GRIS, line=LIGNE, line_w=0.75, adj=[0.12])
    D.add_text_runs(s, 0.82, 4.34, 8.37, 0.42, [
        ([('Principe transversal : ', dict(size=9, bold=True, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
        ([("« alerter, c'est protéger, pas dénoncer » — devoir d'agir si une situation est à risque.", dict(size=9, bold=False, italic=False, color=SLATE, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.MIDDLE, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    return s


def slide_sequencement(prs, media):
    """Sequencement & lancement — 3 phases a pastilles statut rouge/ambre/vert (slide 14 de la v6)."""
    s = _nouvelle_slide(prs, '06 - Slide vide', garder=())
    D.add_text_runs(s, 0.62, 0.34, 7.7, 0.5, [
        ([('Séquencement & lancement du pilote', dict(size=18, bold=True, italic=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 0.62, 0.86, 8, 0.32, [
        ([('Ce qui vient avant, ce qui avance en parallèle, puis le lancement.', dict(size=10.5, bold=False, italic=False, color=SLATE, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    s.shapes.add_picture(io.BytesIO(media['image4.png']), Inches(8.62), Inches(0.28), Inches(0.62), Inches(0.62))
    D.add_forme(s, 'roundRect', 0.62, 1.3, 2.76, 2.74, fill=BLANC, line=LIGNE, line_w=0.75, adj=[0.06])
    D.add_text_runs(s, 0.8, 1.43, 2.4, 0.28, [
        ([('PRÉREQUIS', dict(size=10, bold=True, italic=False, color=CYAN, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'rect', 0.8, 1.74, 0.42, 0.045, fill=CYAN)
    D.add_forme(s, 'roundRect', 0.8, 1.82, 0.935, 0.24, fill='#C03A2B', adj=[0.5])
    D.add_text_runs(s, 0.8, 1.825, 0.935, 0.23, [
        ([('1 · AVANT', dict(size=8, bold=True, italic=False, color=BLANC, font='Outfit'))], dict(align=PP_ALIGN.CENTER, space_after=2)),
    ], anchor=MSO_ANCHOR.MIDDLE, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 0.8, 2.16, 2.4, 1.76, [
        ([('•  Trancher les 3 arbitrages (board P&D) + rédiger la charte d’écoute', dict(size=8, bold=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=4, line_spacing=1.03)),
        ([('•  Formaliser « mission toxique » : signalement → arbitrage → retrait', dict(size=8, bold=False, color=NAVY, font='Outfit'))], dict(space_after=4, line_spacing=1.03)),
        ([('•  Contact d’alerte provisoire actif (coach → lead RH)', dict(size=8, bold=False, color=NAVY, font='Outfit'))], dict(space_after=4, line_spacing=1.03)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'roundRect', 3.63, 1.3, 2.76, 2.74, fill=BLANC, line=LIGNE, line_w=0.75, adj=[0.06])
    D.add_text_runs(s, 3.81, 1.43, 2.4, 0.28, [
        ([('RÉACTIVER & OUTILLER', dict(size=10, bold=True, italic=False, color=CYAN, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'rect', 3.81, 1.74, 0.42, 0.045, fill=CYAN)
    D.add_forme(s, 'roundRect', 3.81, 1.82, 1.46, 0.24, fill='#C88A14', adj=[0.5])
    D.add_text_runs(s, 3.81, 1.825, 1.46, 0.23, [
        ([('2 · EN PARALLÈLE', dict(size=8, bold=True, italic=False, color=BLANC, font='Outfit'))], dict(align=PP_ALIGN.CENTER, space_after=2)),
    ], anchor=MSO_ANCHOR.MIDDLE, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 3.81, 2.16, 2.4, 1.76, [
        ([('•  Quick win < 1 mois : ateliers de co-dev managers, sans dépendance', dict(size=8, bold=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=4, line_spacing=1.03)),
        ([('•  Réactiver le cercle : élaguer « Café sans solution », garder les 2 méthodes', dict(size=8, bold=False, color=NAVY, font='Outfit'))], dict(space_after=4, line_spacing=1.03)),
        ([('•  Réactiver la clinique (process 6 étapes + garde-fous)', dict(size=8, bold=False, color=NAVY, font='Outfit'))], dict(space_after=4, line_spacing=1.03)),
        ([('•  Lancer le questionnaire évaluation + besoins', dict(size=8, bold=False, color=NAVY, font='Outfit'))], dict(space_after=4, line_spacing=1.03)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'roundRect', 6.63, 1.3, 2.76, 2.74, fill=BLANC, line=LIGNE, line_w=0.75, adj=[0.06])
    D.add_text_runs(s, 6.81, 1.43, 2.4, 0.28, [
        ([('LANCER LE PILOTE', dict(size=10, bold=True, italic=False, color=CYAN, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'rect', 6.81, 1.74, 0.42, 0.045, fill=CYAN)
    D.add_forme(s, 'roundRect', 6.81, 1.82, 1.235, 0.24, fill='#2E8B57', adj=[0.5])
    D.add_text_runs(s, 6.81, 1.825, 1.235, 0.23, [
        ([('3 · LANCEMENT', dict(size=8, bold=True, italic=False, color=BLANC, font='Outfit'))], dict(align=PP_ALIGN.CENTER, space_after=2)),
    ], anchor=MSO_ANCHOR.MIDDLE, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_text_runs(s, 6.81, 2.16, 2.4, 1.76, [
        ([('•  Pilote mentorat mission complexe — la priorité', dict(size=8, bold=False, color=NAVY, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=4, line_spacing=1.03)),
        ([('•  1ʳᵉ vague de cercles + communication (charte)', dict(size=8, bold=False, color=NAVY, font='Outfit'))], dict(space_after=4, line_spacing=1.03)),
        ([('•  Restitution en tribune (17 sept)', dict(size=8, bold=False, color=NAVY, font='Outfit'))], dict(space_after=4, line_spacing=1.03)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    D.add_forme(s, 'roundRect', 0.62, 4.5, 8.77, 0.44, fill=GRIS, line=LIGNE, line_w=0.75, adj=[0.12])
    D.add_text_runs(s, 0.82, 4.5, 8.37, 0.44, [
        ([('Cible : ', dict(size=9, bold=True, color=NAVY, font='Outfit')), ('lancement du questionnaire début septembre, tribune de restitution le 17 septembre 2026 — calendrier détaillé page suivante.', dict(size=9, bold=False, color=SLATE, font='Outfit'))], dict(align=PP_ALIGN.LEFT, space_after=2)),
    ], anchor=MSO_ANCHOR.MIDDLE, wrap=True, autosize=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    return s


def slide_roadmap(prs, media):
    """Roadmap datee — timeline 5 jalons, tribune 17 sept en vert (slide 15 de la v6)."""
    s = _nouvelle_slide(prs, '06 - Slide vide', garder=())
    D.add_text_runs(s, 0.62, 0.4, 9, 0.5, [
        ([('Prochaines étapes — calendrier du pilote', dict(size=18, bold=True, italic=False, color=NAVY))], dict(align=PP_ALIGN.LEFT)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.NONE)
    D.add_text_runs(s, 0.62, 0.92, 9, 0.32, [
        ([('De la conception du questionnaire à la tribune de restitution — été → 17 septembre 2026.', dict(size=10.5, bold=False, italic=False, color=SLATE))], dict(align=PP_ALIGN.LEFT)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.NONE)
    D.add_forme(s, 'rect', 1.35, 2.98, 7.24, 0.045, fill=CYAN)
    D.add_forme(s, 'ellipse', 1.26, 2.91, 0.18, 0.18, fill=CYAN, line=BLANC, line_w=1.5)
    D.add_forme(s, 'rect', 1.344, 2.62, 0.012, 0.29, fill='#C9CEDA')
    D.add_text_runs(s, 0.45, 1.3, 1.8, 1.28, [
        ([('Fin juillet', dict(size=10, bold=True, italic=False, color=CYAN))], dict(align=PP_ALIGN.LEFT, space_after=1)),
        ([('Conception du GForm', dict(size=10, bold=True, italic=False, color=NAVY))], dict(align=PP_ALIGN.LEFT, space_after=1)),
        ([('Rédiger le questionnaire : évaluation de l’existant + recueil des besoins.', dict(size=8, bold=False, italic=False, color=SLATE))], dict(align=PP_ALIGN.LEFT)),
    ], anchor=MSO_ANCHOR.BOTTOM, wrap=True, autosize=MSO_AUTO_SIZE.NONE)
    D.add_forme(s, 'ellipse', 3.07, 2.91, 0.18, 0.18, fill=CYAN, line=BLANC, line_w=1.5)
    D.add_forme(s, 'rect', 3.154, 3.09, 0.012, 0.26, fill='#C9CEDA')
    D.add_text_runs(s, 2.26, 3.4, 1.8, 1.3, [
        ([('Début août', dict(size=10, bold=True, italic=False, color=CYAN))], dict(align=PP_ALIGN.LEFT, space_after=1)),
        ([('Partage & feedback', dict(size=10, bold=True, italic=False, color=NAVY))], dict(align=PP_ALIGN.LEFT, space_after=1)),
        ([('Test interne du formulaire, ajustements avant diffusion.', dict(size=8, bold=False, italic=False, color=SLATE))], dict(align=PP_ALIGN.LEFT)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.NONE)
    D.add_forme(s, 'ellipse', 4.88, 2.91, 0.18, 0.18, fill=CYAN, line=BLANC, line_w=1.5)
    D.add_forme(s, 'rect', 4.964, 2.62, 0.012, 0.29, fill='#C9CEDA')
    D.add_text_runs(s, 4.07, 1.3, 1.8, 1.28, [
        ([('Fin août', dict(size=10, bold=True, italic=False, color=CYAN))], dict(align=PP_ALIGN.LEFT, space_after=1)),
        ([('Partage & échange RH', dict(size=10, bold=True, italic=False, color=NAVY))], dict(align=PP_ALIGN.LEFT, space_after=1)),
        ([('Alignement du dispositif et du questionnaire avec les RH.', dict(size=8, bold=False, italic=False, color=SLATE))], dict(align=PP_ALIGN.LEFT)),
    ], anchor=MSO_ANCHOR.BOTTOM, wrap=True, autosize=MSO_AUTO_SIZE.NONE)
    D.add_forme(s, 'ellipse', 6.69, 2.91, 0.18, 0.18, fill=CYAN, line=BLANC, line_w=1.5)
    D.add_forme(s, 'rect', 6.774, 3.09, 0.012, 0.26, fill='#C9CEDA')
    D.add_text_runs(s, 5.88, 3.4, 1.8, 1.3, [
        ([('1–15 sept', dict(size=10, bold=True, italic=False, color=CYAN))], dict(align=PP_ALIGN.LEFT, space_after=1)),
        ([('Lancement du GForm', dict(size=10, bold=True, italic=False, color=NAVY))], dict(align=PP_ALIGN.LEFT, space_after=1)),
        ([('Diffusion large, période de collecte des réponses (2 sem.).', dict(size=8, bold=False, italic=False, color=SLATE))], dict(align=PP_ALIGN.LEFT)),
    ], anchor=MSO_ANCHOR.TOP, wrap=True, autosize=MSO_AUTO_SIZE.NONE)
    D.add_forme(s, 'ellipse', 8.5, 2.91, 0.18, 0.18, fill='#2E8B57', line=BLANC, line_w=1.5)
    D.add_forme(s, 'rect', 8.584, 2.62, 0.012, 0.29, fill='#C9CEDA')
    D.add_text_runs(s, 7.69, 1.3, 1.8, 1.28, [
        ([('17 sept', dict(size=10, bold=True, italic=False, color='#2E8B57'))], dict(align=PP_ALIGN.LEFT, space_after=1)),
        ([('Tribune — restitution', dict(size=10, bold=True, italic=False, color=NAVY))], dict(align=PP_ALIGN.LEFT, space_after=1)),
        ([('Présentation des résultats et des suites en tribune.', dict(size=8, bold=False, italic=False, color=SLATE))], dict(align=PP_ALIGN.LEFT)),
    ], anchor=MSO_ANCHOR.BOTTOM, wrap=True, autosize=MSO_AUTO_SIZE.NONE)
    D.add_forme(s, 'roundRect', 0.62, 4.95, 8.77, 0.42, fill=GRIS, adj=[0.12])
    D.add_text_runs(s, 0.82, 4.95, 8.4, 0.42, [
        ([('Jalon clé : la tribune du 17 septembre 2026 fixe la date de restitution — tout le séquencement en découle.', dict(size=9, bold=False, italic=False, color=NAVY))], dict(align=PP_ALIGN.LEFT)),
    ], anchor=MSO_ANCHOR.MIDDLE, wrap=True, autosize=MSO_AUTO_SIZE.NONE)
    return s


def build():
    """Reconstruit les 15 slides de la v6 dans un fichier NEUF (v7-genere) :
    v6 ouverte comme template (masters/layouts/thème), média extrait, slides
    reconstruites à zéro. Self-check bloquant (géométrie + débordements de
    texte) puis purge des relations de slide orphelines avant save — règles
    deck binaire (cf. pptx_deck, section helpers durcis). Renvoie la liste
    des problèmes (vide = OK)."""
    media = charger_media(SOURCE_V6)
    prs = Presentation(SOURCE_V6)
    D.clear_slides(prs)
    _neutraliser_fond_cadre(prs)

    slide_couverture(prs, media)
    slide_sommaire(prs, media)

    # === Chapitre 01 — Chantier OHC (constat d'enquête + leviers) ===
    slide_chapitre(prs, "01", "Chantier OHC")
    slide_leviers(prs, media)
    slide_personas(prs, media)
    slide_architecture(prs, media)

    # === Chapitre 02 — Existant & Évaluation ===
    slide_chapitre(prs, "02", "Existant & Évaluation")
    slide_existant(prs, media)
    slide_evaluation(prs, media)

    # === Chapitre 03 — Nouveautés (le mentorat, seule vraie brique neuve) ===
    slide_chapitre(prs, "03", "Nouveautés")
    slide_mentorat(prs, media)

    # === Chapitre 04 — Next steps ===
    slide_chapitre(prs, "04", "Next steps")
    slide_arbitrages(prs, media)
    slide_sequencement(prs, media)
    slide_roadmap(prs, media)

    problemes = (D.verifier_geometrie(prs)
                 + D.verifier_debordements_texte(prs))
    if problemes:
        print(f"SELF-CHECK: {len(problemes)} probleme(s)")
        for p in problemes:
            print(" -", p)
    else:
        print("SELF-CHECK: OK — geometrie propre, aucun debordement de texte")

    purges = D.purger_rels_slides_orphelines(prs)
    if purges:
        print(f"purge: {purges} relation(s) de slide orpheline(s) lachee(s)")
    prs.save(SORTIE)
    print("Ecrit:", SORTIE, f"({len(prs.slides)} slides)")
    return problemes


if __name__ == "__main__":
    sys.exit(1 if build() else 0)
