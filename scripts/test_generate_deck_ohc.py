# -*- coding: utf-8 -*-
"""Tests fonctionnels du générateur du deck OHC « dispositif d'écoute ».

Assertions rejouables sur le .pptx réellement généré (modèle : VSCode3
docs/cadrage-ppt/test_generate_deck.py, idées qualité reprises de VSCode2
tests/test_deck_qualite.py), adaptées au deck binaire reconstruit :

  - STRUCTURE : le fichier généré existe, 15 slides, self-check du build vide ;
    chaque slide de contenu retrouvée par TITRE via D.trouver_slide_par_titre
    — son assertion d'unicité sert aussi de test anti-doublon (leçon deck OHC :
    les matchers approximatifs ont tous été piégés) ; les 4 dividers sur le
    layout natif « 51 - Chapitre [2] » avec leur numéro 01-04.
  - GÉOMÉTRIE : verifier_geometrie == [] et verifier_debordements_texte == []
    sur RELECTURE du fichier (pas seulement sur l'objet en mémoire du build).
  - QUALITÉ : aucune ombre portée explicite (règle dure OCTO : différenciation
    par couleur/bordure, jamais d'ombre) ; aucune police générique posée sur
    les runs (Arial/Calibri… — le deck porte Outfit explicite ou l'héritage du
    layout) ; aucune relation de slide orpheline à la relecture (le filet
    anti-0x80CB4404) ; les images héritées du deck original bien posées ;
    aucun texte gabarit résiduel.
  - RENDU RÉEL : conversion PDF LibreOffice headless + comptage de pages == 15
    — SKIP propre et explicite si LibreOffice est absent de la machine (la
    vérification de référence reste l'ouverture PowerPoint COM + rendu PNG,
    faite en session).

Usage : py scripts/test_generate_deck_ohc.py
"""
import os
import re
import subprocess
import sys
import tempfile

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)

from pptx import Presentation  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: E402
from pptx.oxml.ns import qn  # noqa: E402

import pptx_deck as D  # noqa: E402
import generate_deck_ohc as gen  # noqa: E402

echecs = 0


def check(cond, msg):
    global echecs
    if cond:
        print(f"  ok   {msg}")
    else:
        echecs += 1
        print(f"  FAIL {msg}")


def _soffice_path():
    defaut = r"C:\Program Files\LibreOffice\program\soffice.exe"
    return defaut if os.path.exists(defaut) else "soffice"


def _verifier_rendu_reel(pptx_path, n_slides_attendu, tmp_dir):
    """Convertit en PDF via LibreOffice et compte les pages — un .pptx qui
    parse avec python-pptx peut quand même être refusé/tronqué par un vrai
    moteur de rendu (même principe que le test VSCode2 éponyme). Renvoie
    (None, detail) si LibreOffice est indisponible (skip propre)."""
    try:
        result = subprocess.run(
            [_soffice_path(), "--headless", "--convert-to", "pdf",
             "--outdir", tmp_dir, pptx_path],
            capture_output=True, timeout=120,
        )
    except (FileNotFoundError, subprocess.TimeoutExpired) as e:
        return None, f"LibreOffice indisponible ({e}) — vérification réelle non faite"
    if result.returncode != 0:
        return False, f"LibreOffice a échoué : {result.stderr.decode(errors='replace')[:300]}"
    pdf_path = os.path.join(
        tmp_dir, os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf")
    if not os.path.exists(pdf_path):
        return False, "LibreOffice n'a produit aucun PDF"
    pdf_bytes = open(pdf_path, "rb").read()
    if len(pdf_bytes) < 2000:
        return False, "PDF quasi vide — rendu suspect"
    page_count = len(re.findall(rb"/Type\s*/Page[^s]", pdf_bytes))
    if page_count != n_slides_attendu:
        return False, (f"{page_count} page(s) rendue(s) pour "
                       f"{n_slides_attendu} slide(s) exportée(s)")
    return True, f"{page_count} pages rendues, conforme aux {n_slides_attendu} slides"


# Titres attendus (extraits de la v6 — apostrophes typographiques comprises) :
# slide 1-based -> texte EXACT d'une shape de la slide (unicité exigée).
TITRES_CONTENU = {
    1: "Plans d’actions suite à l’enquête OHC chez P&D",
    2: "Sommaire",
    4: "En synthèse - les 3 leviers d’actions prioritaires suite à l’OHC",
    5: "Les personas de l'enquête OHC — 8 populations",
    6: "Architecture de priorisation — où se situe le dispositif d'écoute",
    8: "Rappel — 3 dispositifs d’écoute RH déjà en place chez OCTO",
    9: "Évaluer l’existant & recueillir les besoins — Google Form",
    11: "Mentorat mission complexe — le format et son cadrage",
    13: "Écoute individuelle — arbitrages recommandés",
    14: "Séquencement & lancement du pilote",
    15: "Prochaines étapes — calendrier du pilote",
}

# Dividers : slide 1-based -> (numéro, titre de chapitre) sur le layout 51.
DIVIDERS = {
    3: ("01", "Chantier OHC"),
    7: ("02", "Existant & Évaluation"),
    10: ("03", "Nouveautés"),
    12: ("04", "Next steps"),
}

# Images héritées du média du deck original : slide 1-based -> nb attendu.
IMAGES = {4: 3, 5: 1, 11: 1, 13: 1, 14: 1}

# Dividers (2026-07-23, reprise du pattern réel VSCode3) : chacun porte
# désormais une VRAIE photo dans son cadre (plus le blob navy vidé de la
# version précédente) — 1 image posée par divider.
IMAGES_DIVIDERS = {3: 1, 7: 1, 10: 1, 12: 1}


def main():
    print("Build :")
    problemes = gen.build()
    out = gen.SORTIE
    check(problemes == [], f"self-check du build vide — {len(problemes)} problème(s)")
    check(os.path.exists(out) and os.path.getsize(out) > 500_000,
          f"fichier .pptx écrit, taille plausible "
          f"({os.path.getsize(out) if os.path.exists(out) else 0} octets)")

    prs = Presentation(out)

    print("Structure :")
    check(len(prs.slides) == 15, f"15 slides — reçu {len(prs.slides)}")
    for num, titre in TITRES_CONTENU.items():
        try:
            idx, _ = D.trouver_slide_par_titre(prs, titre)
            check(idx == num - 1,
                  f"slide {num} : titre unique retrouvé à la bonne position « {titre[:45]}… »")
        except ValueError as e:
            check(False, f"slide {num} : {e}")
    for num, (numero, titre) in DIVIDERS.items():
        s = prs.slides[num - 1]
        check(s.slide_layout.name == "51 - Chapitre [2]",
              f"slide {num} : divider sur le layout natif 51 (reçu {s.slide_layout.name!r})")
        textes = [sh.text_frame.text.strip() for sh in s.shapes
                  if getattr(sh, "has_text_frame", False)]
        check(titre in textes, f"slide {num} : titre de chapitre « {titre} »")
        check(numero in textes, f"slide {num} : numéro de chapitre « {numero} »")

    print("Géométrie (relecture du fichier généré) :")
    geo = D.verifier_geometrie(prs)
    check(geo == [], f"verifier_geometrie == [] — {len(geo)} problème(s)")
    deb = D.verifier_debordements_texte(prs)
    check(deb == [], f"verifier_debordements_texte == [] — {len(deb)} problème(s)")

    print("Qualité :")
    # Ombres : aucune ombre portée explicite (charte OCTO — bordures, pas d'ombres).
    ombres = 0
    for s in prs.slides:
        for sh in s.shapes:
            spPr = getattr(sh._element, "spPr", None)
            if spPr is None:
                continue
            eff = spPr.find(qn("a:effectLst"))
            if eff is not None and eff.find(qn("a:outerShdw")) is not None:
                ombres += 1
    check(ombres == 0, f"aucune ombre portée explicite — {ombres} trouvée(s)")

    # Polices : aucune police générique posée sur un run (l'héritage layout ou
    # Outfit explicite portent la charte ; Arial en buFont de puce est admis).
    generiques = {"arial", "calibri", "times new roman", "segoe ui"}
    trouve = set()
    for s in prs.slides:
        for sh in s.shapes:
            if not getattr(sh, "has_text_frame", False):
                continue
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.name and r.font.name.lower() in generiques:
                        trouve.add(r.font.name)
    check(not trouve,
          "aucune police générique posée sur un run"
          + (f" (trouvé {sorted(trouve)})" if trouve else ""))

    # Relations de slide orphelines : 0 attendu à la relecture (filet
    # anti-corruption 0x80CB4404 — leçon payée deux fois sur ce deck).
    purges = D.purger_rels_slides_orphelines(Presentation(out))
    check(purges == 0, f"aucune relation de slide orpheline — {purges} purgée(s)")

    # Images héritées du deck original, posées là où la v6 les portait.
    for num, attendu in IMAGES.items():
        s = prs.slides[num - 1]
        n = sum(1 for sh in s.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE)
        check(n == attendu, f"slide {num} : {attendu} image(s) posée(s) (reçu {n})")

    # Dividers : chacun porte une vraie photo dans son cadre (pattern VSCode3,
    # 2026-07-23) — plus le blob navy vidé de la version précédente. La photo
    # doit occuper EXACTEMENT les bornes du cadre du layout (round2DiagRect,
    # largeur >= 2in — cf. D.trouver_cadre_layout) : un simple compte d'image
    # ne détecterait pas une photo mal cadrée (mauvaise taille/position).
    cadre_layout = D.trouver_cadre_layout(
        gen._layout(prs, "51 - Chapitre [2]").shapes, "round2DiagRect", largeur_min_in=2.0)
    check(cadre_layout is not None, "cadre photo du layout Chapitre retrouvé")
    for num, attendu in IMAGES_DIVIDERS.items():
        s = prs.slides[num - 1]
        images = [sh for sh in s.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
        check(len(images) == attendu,
              f"slide {num} : {attendu} photo(s) de chapitre posée(s) (reçu {len(images)})")
        if images and cadre_layout is not None:
            pic = images[0]
            cl, ct, cw, ch, _, _ = cadre_layout
            check(pic.left == cl and pic.top == ct and pic.width == cw and pic.height == ch,
                  f"slide {num} : photo calée exactement sur le cadre du layout "
                  f"(reçu l={pic.left} t={pic.top} w={pic.width} h={pic.height} ; "
                  f"attendu l={cl} t={ct} w={cw} h={ch})")
            geom = pic._element.spPr.find(qn("a:prstGeom"))
            check(geom is not None and geom.get("prst") == "round2DiagRect",
                  f"slide {num} : photo clippée au prstGeom du cadre (round2DiagRect)")

    # Numéro de chapitre dans le placeholder natif idx1 (2026-07-23, reprise
    # explicite du pattern VSCode2/VSCode3) : la pastille-pilule du template
    # (fond blanc, contour navy) porte le numéro sur UNE seule ligne — pas de
    # gros chiffre dessiné par-dessus la photo (l'ancien scrim de protection
    # n'a plus lieu d'être, le numéro ne touche plus la photo). Le wrap
    # ("01" sur 2 lignes dans le petit encart, piège documenté VSCode2/
    # VSCode3) est le risque réel à couvrir ici.
    for num, (numero, _titre) in DIVIDERS.items():
        s = prs.slides[num - 1]
        ph_num = next((ph for ph in s.placeholders if ph.placeholder_format.idx == 1), None)
        check(ph_num is not None, f"slide {num} : placeholder numéro (idx 1) présent")
        if ph_num is not None:
            paras = [p for p in ph_num.text_frame.paragraphs if p.text.strip()]
            check(len(paras) == 1,
                  f"slide {num} : numéro sur une seule ligne, pas de wrap "
                  f"(reçu {len(paras)} paragraphe(s) non vide(s))")
            texte = ph_num.text_frame.text.strip()
            check(texte == numero,
                  f"slide {num} : placeholder numéro contient « {numero} » (reçu {texte!r})")

    # Aucun texte gabarit résiduel (cadres photo du template).
    texte_complet = "\n".join(
        sh.text_frame.text for s in prs.slides for sh in s.shapes
        if getattr(sh, "has_text_frame", False))
    check("ici mettre une Photo" not in texte_complet,
          "aucun texte gabarit de cadre photo résiduel")

    print("Rendu réel (LibreOffice — conversion PDF, comptage de pages) :")
    with tempfile.TemporaryDirectory(prefix="test-deck-ohc-") as tmp:
        ok, detail = _verifier_rendu_reel(out, len(prs.slides), tmp)
        if ok is None:
            print(f"  SKIP {detail}")
        else:
            check(ok, detail)

    print("\nTOUS LES TESTS PASSENT" if echecs == 0 else f"\n{echecs} TEST(S) EN ECHEC")
    sys.exit(0 if echecs == 0 else 1)


if __name__ == "__main__":
    main()
