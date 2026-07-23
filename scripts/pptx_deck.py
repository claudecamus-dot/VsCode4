"""pptx_deck — bibliotheque d'aide pour construire des slides python-pptx
"de qualite" : echelle typographique coherente, formes (barres, jauge, cartes),
couleurs, et surtout un controle geometrique automatique (`verifier_geometrie`)
qui detecte toute forme qui sort de la slide — le defaut classique des decks
generes a la main.

Reprise VSCode4 (arbitrage 2026-07-23, diagnostic superviseur) : copie du module
de reference VSCode2 `app/services/pptx_deck.py` (sur-ensemble de la version
VSCode3), completee d'une section « helpers durcis deck binaire » qui fige en
code les lecons payees sur le deck OHC (runs du 2026-07-21) :
 - recherche de slide par TITRE avec assertion d'unicite (les matchers par
   position, par title_of ou par corps de texte ont tous ete pieges) ;
 - suppression de slide avec drop_rel (sans lui : parts orphelines -> PowerPoint
   refuse d'ouvrir, HRESULT 0x80CB4404) + purge de rels orphelines avant save ;
 - regle AJOUTER-AVANT-SUPPRIMER quand on remplace une slide (noms de parts
   frais, jamais de reutilisation d'un nom de part supprime dans le meme cycle).

Reutilisable hors de ce projet : aucune dependance au domaine metier ici.
Les coordonnees des helpers sont exprimees en POUCES (float) pour la lisibilite.
"""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

# --- Echelle typographique (pt) — une seule source de verite ---
# `title` = 20 (aligne sur la charte de reference bmad-iap-cadrage-synthese,
# dont les titres de contenu sont a ~18-20pt — 26 faisait des titres
# surdimensionnes qui debordaient/wrappaient).
TYPE = {
    "title": 20, "h2": 18, "h3": 14, "body": 12, "small": 10.5, "tiny": 9,
    "kpi": 44, "kpi_unit": 16,
}

PALETTE = ["#2c5cc5", "#1e6b34", "#b3261e", "#b8860b", "#6a3d9a", "#138086"]

INK = "#1c2330"
MUTED = "#6b7280"
LINE = "#e6e8ee"
TRACK = "#eef1f7"
OK = "#1e6b34"
WARN = "#b3261e"
GOLD = "#b8860b"


def rgb(hexa):
    return RGBColor.from_string(hexa.lstrip("#"))


def melanger_blanc(hexa, frac):
    """Éclaircit une couleur en la mélangeant avec du blanc. `frac`=0.0 -> couleur
    d'origine, 1.0 -> blanc pur. Sert aux fonds teintés (cellules de matrice,
    encarts) qui doivent rester lisibles sous du texte foncé."""
    frac = max(0.0, min(1.0, float(frac)))
    h = hexa.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    r = round(r + (255 - r) * frac)
    g = round(g + (255 - g) * frac)
    b = round(b + (255 - b) * frac)
    return f"#{r:02x}{g:02x}{b:02x}"


def couleur_pilier(i):
    return PALETTE[i % len(PALETTE)]


def _no_shadow(shape):
    # Les autoshapes heritent parfois d'une ombre du theme : on la coupe.
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


# --- Police effective du deck ---
# Deux sources possibles pour la police d'un template OCTO :
#  - police_marque() : la famille dominante sur les PLACEHOLDERS des layouts
#    (Outfit sur OCTO). C'est la charte "de conception" — mais Outfit n'est pas
#    une police systeme : sur une machine ou elle n'est pas installee, PowerPoint
#    la rend en SUBSTITUTION (rendu visuellement != charte).
#  - police_theme() : le fontScheme du theme (minorFont = corps). Sur OCTO c'est
#    Arial, police systeme => rendu garanti.
# set_police() pose la famille ; add_text l'applique a chaque run dessine, et
# appliquer_police() la force sur les runs de PLACEHOLDER (titres/couverture/
# chapitre, qui heriteraient sinon l'Outfit du layout). POLICE = None -> heritage.
POLICE = None

_SUFFIXES_POIDS = (" SemiBold", " Semibold", " Semi Bold", " ExtraBold",
                   " Extra Bold", " Bold", " Medium", " Light", " Regular",
                   " Black", " Thin", " ExtraLight", " Extra Light")


def set_police(nom):
    global POLICE
    POLICE = nom or None


def _famille_police(typeface):
    for suf in _SUFFIXES_POIDS:
        if typeface.endswith(suf):
            return typeface[: -len(suf)].strip()
    return typeface


def police_marque(prs):
    """Detecte la police de marque du template : la famille (suffixe de poids
    retire) la plus frequente sur les placeholders des layouts. Placeholders et
    pas fontScheme : sur les templates OCTO, les titres portent la charte (Outfit,
    decline en poids nommes) alors que le fontScheme peut n'etre qu'un repli
    generique (Arial). Les references de theme (+mn-lt/+mj-lt) sont ignorees.
    Renvoie None si rien d'exploitable (l'appelant garde l'heritage par defaut)."""
    import re
    from collections import Counter
    try:
        layouts = prs.slide_masters[0].slide_layouts
    except Exception:
        return None
    compte = Counter()
    for lay in layouts:
        for ph in lay.placeholders:
            for tf in re.findall(r'<a:latin typeface="([^"]+)"', ph._element.xml):
                if tf.startswith("+"):
                    continue
                compte[_famille_police(tf)] += 1
    return compte.most_common(1)[0][0] if compte else None


def police_theme(prs):
    """Police du THEME du template (fontScheme, minorFont = corps de texte).
    Sur OCTO c'est Arial — une police systeme, donc rendue telle quelle par
    PowerPoint, contrairement a l'Outfit des placeholders (non installee ->
    substitution). Renvoie la famille (suffixe de poids retire), ou None."""
    import re
    xml = None
    try:
        xml = prs.slide_masters[0].part.part_related_by(RT.THEME).blob.decode(
            "utf-8", "ignore")
    except Exception:
        xml = None
    if xml is None:
        try:
            for part in prs.part.package.iter_parts():
                if "theme" in str(part.partname):
                    xml = part.blob.decode("utf-8", "ignore")
                    break
        except Exception:
            return None
    if not xml:
        return None
    m = re.search(r'<a:minorFont>.*?<a:latin typeface="([^"]*)"', xml, re.S)
    if m and m.group(1) and not m.group(1).startswith("+"):
        return _famille_police(m.group(1))
    return None


def appliquer_police(text_frame):
    """Force la police effective du deck (POLICE) sur tous les runs d'un
    text_frame de PLACEHOLDER. Les placeholders (titre natif, couverture,
    intercalaire de chapitre) heritent la police du layout — l'Outfit non
    installee des templates OCTO — qui serait rendue en substitution ; on la
    remplace par POLICE pour que le texte de placeholder coincide avec le texte
    dessine par add_text. No-op si POLICE est None."""
    if not POLICE:
        return
    for p in text_frame.paragraphs:
        for run in p.runs:
            run.font.name = POLICE


def add_text(slide, l, t, w, h, lignes, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT,
             wrap=True):
    """Ajoute une zone de texte. `lignes` = liste de (texte, opts) ; chaque
    element devient un paragraphe. opts: size,bold,italic,color,align,
    space_before,space_after,line_spacing."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    for i, (texte, opts) in enumerate(lignes):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get("align", align)
        if "space_before" in opts:
            p.space_before = Pt(opts["space_before"])
        if "space_after" in opts:
            p.space_after = Pt(opts["space_after"])
        if "line_spacing" in opts:
            p.line_spacing = opts["line_spacing"]
        run = p.add_run()
        run.text = texte
        f = run.font
        f.size = Pt(opts.get("size", TYPE["body"]))
        f.bold = opts.get("bold", False)
        f.italic = opts.get("italic", False)
        f.color.rgb = rgb(opts.get("color", INK))
        nom = opts.get("font", POLICE)  # police de marque du deck (cf. set_police)
        if nom:
            f.name = nom
    return box


def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=1.0, rounded=False,
             radius=0.12):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    _no_shadow(shp)
    if rounded:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = rgb(line)
        shp.line.width = Pt(line_w)
    shp.text_frame.paragraphs[0].text = ""
    return shp


# --- Formes et texte « riches » (ajout 2026-07-23, générateur deck OHC) ---
# Nés de la reconstruction fidèle d'un deck binaire : add_text force taille et
# couleur sur chaque run (parfait pour dessiner du neuf, faux pour reproduire
# du texte de placeholder qui doit HÉRITER sa charte du layout), et add_rect ne
# couvre que les rectangles. Les helpers ci-dessous ne posent QUE les
# propriétés fournies.

FORMES_PRST = {
    "rect": MSO_SHAPE.RECTANGLE,
    "roundRect": MSO_SHAPE.ROUNDED_RECTANGLE,
    "round1Rect": MSO_SHAPE.ROUND_1_RECTANGLE,
    "round2DiagRect": MSO_SHAPE.ROUND_2_DIAG_RECTANGLE,
    "round2SameRect": MSO_SHAPE.ROUND_2_SAME_RECTANGLE,
    "ellipse": MSO_SHAPE.OVAL,
    "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "rtTriangle": MSO_SHAPE.RIGHT_TRIANGLE,
    "diamond": MSO_SHAPE.DIAMOND,
    "tear": MSO_SHAPE.TEAR,
}


def add_forme(slide, prst, l, t, w, h, fill=None, line=None, line_w=1.0,
              adj=None, rot=0, dash=None, fill_alpha=None):
    """Autoshape générique par nom de preset OOXML (« roundRect »,
    « round2DiagRect », « ellipse », « triangle », « tear »… cf. FORMES_PRST).
    Complète add_rect (limité aux rectangles) pour reproduire fidèlement les
    formes d'un deck existant : `adj` = liste d'adjustments posés dans l'ordre
    du preset (fractions de l'échelle OOXML 100000, ex. 0.16667 pour
    `val 16667`), `rot` en degrés, `dash` = style de tirets OOXML de la
    bordure (« dot », « dash »… — None = trait plein). Ombre coupée (règle
    dure OCTO), fill/line hexa ou None (fond transparent / sans bordure).
    `fill_alpha` (0-100, None = opaque) pose une opacité PARTIELLE sur `fill` —
    un scrim plat semi-transparent (ex. lisibilité d'un texte posé sur une
    photo) reste conforme à la charte « pas d'ombre portée » puisqu'il n'a
    aucun flou/décalage, contrairement à `a:outerShdw` (interdit ailleurs
    dans ce module via `_no_shadow`)."""
    shp = slide.shapes.add_shape(FORMES_PRST[prst], Inches(l), Inches(t),
                                 Inches(w), Inches(h))
    _no_shadow(shp)
    if adj:
        for i, v in enumerate(adj):
            try:
                shp.adjustments[i] = v
            except Exception:
                pass
    if rot:
        shp.rotation = rot
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(fill)
        if fill_alpha is not None:
            srgb = shp.fill.fore_color._xFill.find(qn("a:srgbClr"))
            if srgb is not None:
                for old in srgb.findall(qn("a:alpha")):
                    srgb.remove(old)
                a = srgb.makeelement(
                    qn("a:alpha"), {"val": str(int(round(fill_alpha * 1000)))})
                srgb.append(a)
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = rgb(line)
        shp.line.width = Pt(line_w)
        if dash:
            ln = shp.line._get_or_add_ln()
            pd = ln.find(qn("a:prstDash"))
            if pd is None:
                pd = ln.makeelement(qn("a:prstDash"), {})
                ln.append(pd)
            pd.set("val", dash)
    shp.text_frame.paragraphs[0].text = ""
    return shp


def definir_geometrie(shape, l, t, w, h):
    """Pose la géométrie (pouces) d'une shape existante — typiquement un
    placeholder cloné du layout dont on veut figer la position telle que
    mesurée sur le deck de référence (un placeholder sans xfrm explicite
    hérite du layout : expliciter la même valeur est sans effet visuel mais
    rend la reproduction indépendante d'une évolution du layout)."""
    shape.left = Inches(l)
    shape.top = Inches(t)
    shape.width = Inches(w)
    shape.height = Inches(h)


def configurer_text_frame(tf, anchor=None, wrap=None, autosize=None,
                          margins=None):
    """Configure un text frame en ne touchant QUE les propriétés fournies
    (None = laisser hériter) : ancrage vertical, retour à la ligne, auto-size
    (MSO_AUTO_SIZE) et marges internes `(gauche, haut, droite, bas)` en
    pouces. Indispensable pour reproduire un deck existant sans écraser les
    réglages hérités des layouts."""
    if anchor is not None:
        tf.vertical_anchor = anchor
    if wrap is not None:
        tf.word_wrap = wrap
    if autosize is not None:
        tf.auto_size = autosize
    if margins is not None:
        ml, mt, mr, mb = margins
        tf.margin_left = Inches(ml)
        tf.margin_top = Inches(mt)
        tf.margin_right = Inches(mr)
        tf.margin_bottom = Inches(mb)


def definir_paragraphes(tf, paras):
    """Écrit des paragraphes « riches » (plusieurs runs stylés PAR paragraphe,
    ex. un mot en gras au milieu d'une phrase) en ne posant que les propriétés
    fournies — contrairement à add_text qui force taille et couleur sur chaque
    run, ce qui casserait l'héritage de charte du texte de placeholder.
    `paras` = liste de (runs, opts_para) ; `runs` = liste de (texte, opts_run).
    opts_run : size (pt), bold, italic, color (hexa), font ; opts_para :
    align (PP_ALIGN), space_before/space_after (pt), line_spacing (multiple),
    bullet (dict char/size/font/color — puce réelle buChar avec retrait
    suspendu marL/indent en pouces, defaut 0.1875), marL/indent (pouces).
    Le contenu existant du text frame est remplacé."""
    tf.clear()
    for i, (runs, opts) in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if "align" in opts:
            p.alignment = opts["align"]
        if "space_before" in opts:
            p.space_before = Pt(opts["space_before"])
        if "space_after" in opts:
            p.space_after = Pt(opts["space_after"])
        if "line_spacing" in opts:
            p.line_spacing = opts["line_spacing"]
        if opts.get("bullet") or "marL" in opts or "indent" in opts:
            bu = opts.get("bullet") or {}
            marL = opts.get("marL", 0.1875 if bu else 0.0)
            indent = opts.get("indent", -marL if bu else 0.0)
            pPr = p._p.get_or_add_pPr()
            pPr.set("marL", str(int(Inches(marL))))
            pPr.set("indent", str(int(Inches(indent))))
            if bu:
                for tag, attrs in (
                        ("a:buClr", None),
                        ("a:buSzPts", {"val": str(int(bu.get("size", 10) * 100))}),
                        ("a:buFont", {"typeface": bu.get("font", "Arial")}),
                        ("a:buChar", {"char": bu.get("char", "•")})):
                    el = pPr.makeelement(qn(tag), attrs or {})
                    if tag == "a:buClr":
                        clr = pPr.makeelement(qn("a:srgbClr"), {
                            "val": bu.get("color", INK).lstrip("#").upper()})
                        el.append(clr)
                    pPr.append(el)
        for texte, ro in runs:
            r = p.add_run()
            r.text = texte
            f = r.font
            if ro.get("size") is not None:
                f.size = Pt(ro["size"])
            if ro.get("bold") is not None:
                f.bold = ro["bold"]
            if ro.get("italic") is not None:
                f.italic = ro["italic"]
            if ro.get("color"):
                f.color.rgb = rgb(ro["color"])
            if ro.get("font"):
                f.name = ro["font"]


def add_text_runs(slide, l, t, w, h, paras, anchor=None, wrap=True,
                  autosize=None, margins=(0, 0, 0, 0)):
    """Zone de texte « riche » : la version multi-runs d'add_text (un
    paragraphe peut mélanger des runs de styles différents), configurée via
    configurer_text_frame + definir_paragraphes. Mêmes conventions d'unités
    (pouces / pt)."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    configurer_text_frame(box.text_frame, anchor=anchor, wrap=wrap,
                          autosize=autosize, margins=margins)
    definir_paragraphes(box.text_frame, paras)
    return box


def add_hbar(slide, l, t, w, h, frac, fill, track=TRACK):
    """Barre de progression horizontale (piste + remplissage), coins arrondis."""
    frac = max(0.0, min(1.0, float(frac)))
    add_rect(slide, l, t, w, h, fill=track, rounded=True, radius=0.5)
    if frac > 0:
        wv = max(h, w * frac)  # largeur minimale visible = hauteur (pastille)
        add_rect(slide, l, t, wv, h, fill=fill, rounded=True, radius=0.5)


def add_gauge(slide, l, t, size, frac, fill, track=TRACK, hole=62):
    """Jauge circulaire (anneau) via un graphique doughnut a 2 segments.
    Renvoie le GraphicFrame. Le libelle central est a poser separement."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    frac = max(0.0, min(1.0, float(frac)))
    data = CategoryChartData()
    data.categories = ["v", "r"]
    data.add_series("g", (frac, 1 - frac))
    gf = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(l), Inches(t),
                                Inches(size), Inches(size), data)
    chart = gf.chart
    chart.has_legend = False
    chart.has_title = False
    plot = chart.plots[0]
    plot.has_data_labels = False
    # Taille du trou
    dough = plot._element  # <c:doughnutChart>
    hs = dough.find(qn("c:holeSize"))
    if hs is None:
        hs = dough.makeelement(qn("c:holeSize"), {"val": str(hole)})
        dough.append(hs)
    else:
        hs.set("val", str(hole))
    # Couleurs des 2 segments
    pts = plot.series[0].points
    for pt_, col in ((pts[0], fill), (pts[1], track)):
        pt_.format.fill.solid()
        pt_.format.fill.fore_color.rgb = rgb(col)
        pt_.format.line.color.rgb = rgb("#ffffff")
        pt_.format.line.width = Pt(1)
    return gf


def add_card(slide, l, t, w, h, accent=None):
    """Carte blanche a coins arrondis, fine bordure grise, + lisere couleur a
    gauche. `accent` None -> pas de lisere (repli)."""
    add_rect(slide, l, t, w, h, fill="#ffffff", line=LINE, line_w=0.75,
             rounded=True, radius=0.06)
    if accent:
        add_rect(slide, l, t, 0.07, h, fill=accent, rounded=True, radius=0.5)


def add_card_header(slide, l, t, w, label, color, size=None):
    """En-tete de carte facon OCTO : libelle en petites capitales couleur
    `color` + court filet d'accent dessous. Retourne le y ou le contenu peut demarrer."""
    size = TYPE["h3"] if size is None else size
    pad = 0.0
    add_text(slide, l, t, w, 0.34,
             [(label.upper(), dict(size=size, bold=True, color=color))])
    add_rect(slide, l + pad, t + 0.36, 0.5, 0.045, fill=color)
    return t + 0.36 + 0.045 + 0.14  # y de depart du contenu sous le filet


def add_dot(slide, x, y, d, color):
    """Petite pastille ronde pleine (puce de legende / marqueur de chip)."""
    return add_rect(slide, x, y, d, d, fill=color, rounded=True, radius=0.5)


def add_chip(slide, x, y, w, h, label, color, text_color="#ffffff", size=None,
             outline=False):
    """Pastille etiquette (pill) a coins pleins arrondis — tag de categorisation
    (« Quick win », « DECISION », un numero de rang…). `outline=True` : fond
    blanc, bordure + texte de la couleur (variante sobre) ; sinon fond plein
    `color`, texte `text_color`. Texte centre, sans ombre."""
    size = TYPE["tiny"] if size is None else size
    if outline:
        add_rect(slide, x, y, w, h, fill="#ffffff", line=color, line_w=1.0,
                 rounded=True, radius=0.5)
        txt = color
    else:
        add_rect(slide, x, y, w, h, fill=color, rounded=True, radius=0.5)
        txt = text_color
    add_text(slide, x, y, w, h,
             [(label, dict(size=size, bold=True, color=txt, align=PP_ALIGN.CENTER))],
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)


def add_badge(slide, x, y, d, glyph, color, text_color="#ffffff", size=None,
              bold=True, radius=0.28):
    """Badge-icone : petite tuile carree a coins arrondis (couleur pleine, sans
    ombre — regle dure OCTO : differenciation par couleur/bordure, pas d'ombre)
    portant un glyphe monochrome ou un numero centre. `radius` : 0.5 = pastille
    ronde, ~0.28 = tuile. `bold` a False pour un glyphe qui « tofu » en gras
    dans la police du template. Retourne la tuile."""
    size = TYPE["h3"] if size is None else size
    shp = add_rect(slide, x, y, d, d, fill=color, rounded=True, radius=radius)
    add_text(slide, x, y, d, d,
             [(glyph, dict(size=size, bold=bold, color=text_color, align=PP_ALIGN.CENTER))],
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    return shp


def add_teardrop(slide, x, y, d, label, color, size=None, rot=180, line_w=1.75):
    """Badge goutte (prstGeom « tear ») a CONTOUR — signature OCTO du sommaire :
    fond blanc, bordure couleur, numero/label centre (pose en zone de texte
    separee pour rester droit quand la goutte est pivotee). `rot` oriente la
    pointe (defaut 180 = pointe en bas-gauche). Sans ombre (regle dure)."""
    size = TYPE["h2"] if size is None else size
    shp = slide.shapes.add_shape(MSO_SHAPE.TEAR, Inches(x), Inches(y), Inches(d), Inches(d))
    _no_shadow(shp)
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb("#ffffff")
    shp.line.color.rgb = rgb(color)
    shp.line.width = Pt(line_w)
    shp.rotation = rot
    shp.text_frame.paragraphs[0].text = ""
    add_text(slide, x, y, d, d,
             [(label, dict(size=size, bold=True, color=color, align=PP_ALIGN.CENTER))],
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    return shp


ENCART_BG = "#eceef2"  # gris clair d'encart (structural, cf. LINE/TRACK) — motif OCTO reel


def add_encart(slide, l, t, w, h, text, accent=None, label=None, size=None,
               align=PP_ALIGN.CENTER):
    """Encart « a retenir / so-what » — boite GRIS CLAIR arrondie, sans ombre,
    texte fonce. Un encart est QUIET (gris), pas une bande de couleur criarde :
    la couleur est un accent, pas de la decoration. `accent` pose un fin lisere
    gauche colore ; `label` un prefixe gras. Composant unique reutilise partout."""
    size = TYPE["h3"] if size is None else size
    add_rect(slide, l, t, w, h, fill=ENCART_BG, rounded=True, radius=0.12)
    pad = 0.24
    if accent:
        add_rect(slide, l, t, 0.06, h, fill=accent, rounded=True, radius=0.5)
        pad = 0.28
    lignes = []
    if label:
        lignes.append((label, dict(size=size, bold=True, color=INK, align=align)))
    lignes.append((text, dict(size=size, bold=(label is None), color=INK, align=align)))
    add_text(slide, l + pad, t, w - 2 * pad, h, lignes, anchor=MSO_ANCHOR.MIDDLE, align=align)


def add_range_bar(slide, l, t, w, h, mn, mx, scale_max, fill, marker=None,
                  track=TRACK):
    """Barre d'amplitude min..max sur une echelle 0..scale_max (piste complete +
    segment colore couvrant la plage). `marker` (ex. moyenne) pose un repere
    vertical. Sert a montrer une dispersion sur l'echelle reelle, pas en relatif."""
    add_rect(slide, l, t, w, h, fill=track, rounded=True, radius=0.5)
    fa = max(0.0, min(1.0, mn / scale_max))
    fb = max(0.0, min(1.0, mx / scale_max))
    seg_w = max(h, w * (fb - fa))  # largeur minimale = hauteur (pastille)
    add_rect(slide, l + w * fa, t, seg_w, h, fill=fill, rounded=True, radius=0.5)
    if marker is not None:
        fm = max(0.0, min(1.0, marker / scale_max))
        mx_x = l + w * fm - 0.015
        add_rect(slide, mx_x, t - 0.05, 0.03, h + 0.10, fill=INK, rounded=True,
                 radius=0.5)


def _compte_lignes(texte, cpl):
    """Nombre de lignes apres un repli mot-a-mot pour une largeur de `cpl`
    caracteres. Estimateur volontairement simple (pas de mesure de police reelle) :
    calibre empiriquement via `cpl`, pas cense etre pixel-parfait."""
    total = 0
    for para in str(texte).split("\n"):
        cur = 0
        n = 1
        for mot in para.split():
            ajout = (1 if cur else 0) + len(mot)
            if cur + ajout > cpl and cur:
                n += 1
                cur = len(mot)
            else:
                cur += ajout
        total += n
    return total


def estimer_lignes(texte, largeur_in, taille_pt, cpi_ref=11.0, taille_ref=10.5):
    """Estime le nombre de lignes qu'occupera `texte` une fois reparti mot-a-mot
    sur `largeur_in` pouces a la taille de police `taille_pt`. Les caracteres par
    pouce sont derives de `cpi_ref` (calibre a `taille_ref` pt) par une regle de
    trois : une police 2x plus petite loge environ 2x plus de caracteres/pouce."""
    if not texte:
        return 1
    cpi = cpi_ref * (taille_ref / taille_pt)
    cpl = max(6, int(largeur_in * cpi))
    return _compte_lignes(texte, cpl)


def ajuster_police(textes, largeur_in, taille_max, taille_min, budget_ok, pas=0.5,
                   cpi_ref=11.0, taille_ref=10.5):
    """Adapte la taille de police a la longueur des phrases a restituer : cherche,
    par pas de `pas` pt entre `taille_max` et `taille_min`, la plus GRANDE taille
    telle que `budget_ok(taille, lignes_max)` soit vrai — ou `lignes_max` est le
    nombre de lignes necessaires au plus long de `textes` une fois reparti sur
    `largeur_in` pouces a cette taille (voir `estimer_lignes`).

    `budget_ok` encapsule la contrainte geometrique propre a l'appelant (ex. :
    n cartes empilees doivent tenir dans la bande disponible) — cette fonction
    reste agnostique du domaine. Objectif : ne JAMAIS tronquer/faire deborder une
    phrase — si aucune taille ne satisfait `budget_ok`, on degrade sur
    `taille_min` (texte tres dense) plutot que de laisser un texte coupe.

    Renvoie (taille, lignes_max)."""
    taille = taille_max
    while True:
        lignes_max = max((estimer_lignes(t, largeur_in, taille, cpi_ref, taille_ref)
                          for t in textes), default=1)
        if budget_ok(taille, lignes_max) or taille <= taille_min:
            return (max(taille, taille_min), lignes_max)
        taille = max(taille_min, round(taille - pas, 2))


def tronquer_a_lignes(texte, largeur_in, taille_pt, max_lignes, cpi_ref=11.0,
                      taille_ref=10.5):
    """Tronque `texte` (avec une ellipse finale) pour qu'il tienne dans
    `max_lignes` lignes une fois reparti sur `largeur_in` pouces a `taille_pt`.
    Dernier recours quand meme `ajuster_police` a sa taille plancher ne suffit
    plus a eviter un debordement geometrique — mieux vaut un texte coupe
    proprement qu'une forme qui deborde de la slide. Ne fait rien si le texte
    tient deja dans `max_lignes`."""
    if estimer_lignes(texte, largeur_in, taille_pt, cpi_ref, taille_ref) <= max_lignes:
        return texte
    cpi = cpi_ref * (taille_ref / taille_pt)
    cpl = max(6, int(largeur_in * cpi))
    limite = max(1, cpl * max_lignes - 1)
    tronque = str(texte)[:limite].rstrip()
    dernier_espace = tronque.rfind(" ")
    if dernier_espace > limite * 0.6:
        tronque = tronque[:dernier_espace]
    return tronque.rstrip(" ,;:.") + "…"


def verifier_debordements_texte(prs, cpi_pessimiste=10.7, tolerance_in=0.15):
    """Filet « le texte tient dans sa boite » — complementaire de
    verifier_geometrie (qui ne voit que les BORDS des formes, pas le rendu du
    texte dedans). Pour chaque zone de texte dessinee (wrap actif, ancrage TOP,
    auto-size NONE), estime la hauteur du contenu avec une calibration
    PESSIMISTE (`cpi_pessimiste` < la calibration nominale de estimer_lignes) et
    signale les boites dont le contenu estime depasse la hauteur + tolerance.
    Renvoie une liste de constats (vide = OK) ; l'appelant decide (test dur, ou log)."""
    problemes = []
    for num, slide in enumerate(prs.slides, start=1):
        for sh in slide.shapes:
            if not getattr(sh, "has_text_frame", False):
                continue
            if getattr(sh, "is_placeholder", False):
                continue  # placeholders (titres, couverture…) : PowerPoint les
                # laisse grandir sans cadre visuel — pas le defaut chasse ici
            tf = sh.text_frame
            try:
                if not tf.word_wrap or tf.auto_size != MSO_AUTO_SIZE.NONE:
                    continue
                if tf.vertical_anchor not in (None, MSO_ANCHOR.TOP):
                    continue  # MIDDLE/BOTTOM : contenu deja borne par l'appelant
                if getattr(sh, "rotation", 0):
                    continue  # labels rotes : geometrie non comparable
                w_in = Emu(sh.width).inches
                h_in = Emu(sh.height).inches
            except Exception:
                continue
            if w_in <= 0 or h_in <= 0:
                continue
            est = 0.0
            texte_court = ""
            for p in tf.paragraphs:
                t = "".join(r.text for r in p.runs)
                if not t.strip():
                    continue
                # max des runs styles, pas le premier (un prefixe petit devant
                # un corps plus grand sous-estimerait toute la hauteur)
                tailles = [r.font.size.pt for r in p.runs if r.font.size]
                taille = max(tailles) if tailles else 12.0
                lignes = estimer_lignes(t, w_in, taille, cpi_ref=cpi_pessimiste)
                # meme modele de hauteur de ligne que le layout (le +4/72 couvre
                # deja le space_after usuel — pas de double comptage)
                est += lignes * (taille * 0.017 + 4 / 72)
                texte_court = texte_court or t[:40]
            if est > h_in + tolerance_in:
                problemes.append(
                    f"slide {num}: texte ~{est:.2f}in > boîte {h_in:.2f}in "
                    f"(« {texte_court}… »)"
                )
    return problemes


def paginer_items(items, hauteur_fn, capacite_in):
    """Repartit gloutonnement `items` (ordre preserve) en pages telles que,
    pour chaque page, la somme de `hauteur_fn(item)` (pouces) ne depasse pas
    `capacite_in`. Un item dont la hauteur depasse `capacite_in` a lui seul
    reste seul sur sa page plutot que d'etre coupe — garantit qu'au moins un
    item est consomme par page, donc une terminaison en au plus `len(items)`
    pages. Renvoie une liste de listes (au moins une page, eventuellement
    vide, si `items` est vide)."""
    if not items:
        return [[]]
    pages: list[list] = []
    current: list = []
    current_h = 0.0
    for item in items:
        h = hauteur_fn(item)
        if current and current_h + h > capacite_in:
            pages.append(current)
            current, current_h = [], 0.0
        current.append(item)
        current_h += h
    pages.append(current)
    return pages


def verifier_geometrie(prs, marge_in=0.02):
    """Retourne la liste des problemes : toute forme dont les bords depassent la
    slide (au-dela d'une petite marge de tolerance). Liste vide = OK."""
    W, H = prs.slide_width, prs.slide_height
    tol = Inches(marge_in)
    problemes = []
    for si, slide in enumerate(prs.slides, start=1):
        for shp in slide.shapes:
            try:
                l, t, w, h = shp.left, shp.top, shp.width, shp.height
            except Exception:
                continue
            if None in (l, t, w, h):
                continue
            nom = shp.name or "shape"
            if l < -tol or t < -tol or (l + w) > W + tol or (t + h) > H + tol:
                problemes.append(
                    f"slide {si}: '{nom}' hors cadre "
                    f"(l={Emu(l).inches:.2f} t={Emu(t).inches:.2f} "
                    f"r={Emu(l + w).inches:.2f} b={Emu(t + h).inches:.2f} ; "
                    f"slide {Emu(W).inches:.2f}x{Emu(H).inches:.2f})")
    return problemes


def theme_colors(prs):
    """Lit le nuancier du theme (dk1/lt1/dk2/lt2/accent1..6) du 1er master et le
    renvoie en dict {nom: '#RRGGBB'}. Sert a adapter le deck a la charte du
    template fourni sans rien coder en dur. Renvoie {} si le theme est
    introuvable (l'appelant prevoit un repli)."""
    import re
    try:
        part = prs.slide_masters[0].part
        theme_part = next((r.target_part for r in part.rels.values()
                           if "theme" in r.reltype), None)
        if theme_part is None:
            return {}
        xml = theme_part.blob.decode("utf-8", "ignore")
    except Exception:
        return {}
    m = re.search(r"<a:clrScheme.*?</a:clrScheme>", xml, re.S)
    if not m:
        return {}
    seg = m.group(0)
    out = {}
    for name in ("dk1", "lt1", "dk2", "lt2", "accent1", "accent2",
                 "accent3", "accent4", "accent5", "accent6"):
        mm = re.search(
            r"<a:" + name + r">.*?(?:srgbClr val=\"([0-9A-Fa-f]{6})\"|"
            r"sysClr[^>]*lastClr=\"([0-9A-Fa-f]{6})\")", seg, re.S)
        if mm:
            out[name] = "#" + (mm.group(1) or mm.group(2)).upper()
    return out


# ---------------------------------------------------------------------------
# Cadre photo (frame) NON groupe — complement de pptx-framed-image
# ---------------------------------------------------------------------------
# La skill pptx-framed-image (.claude/skills/pptx-framed-image/scripts/
# framed_image.py, frame_geometry()) suppose que le cadre « ici mettre une
# Photo » est niche dans un GROUP (cas du « cadre blanc » OCTO). Le layout
# « 51 - Chapitre [2] » de CE template pose son cadre comme une AUTOSHAPE de
# premier niveau (pas de group) : frame_geometry() ne s'applique pas telle
# quelle. trouver_cadre_layout() est la variante top-level manquante — assez
# generique pour vivre ici plutot que dans le generateur OHC (constat
# VSCode3 : sa propre variante locale `_find_frame_by_geom` etait deja
# generique, juste jamais remontee dans une bibliotheque partagee).


def trouver_cadre_layout(shapes, prst, largeur_min_in=None):
    """Cherche, parmi les shapes de premier niveau d'un LAYOUT, la premiere
    dont le `<a:prstGeom>` porte le preset `prst` (ex. « round2DiagRect »).
    `largeur_min_in` leve l'ambiguite si plusieurs shapes du layout partagent
    le meme preset a des tailles differentes (regle projet : contraindre un
    matcher par DIMENSIONS, pas seulement par nom/position — une proximite
    seule a deja repeint la mauvaise forme sur ce deck).

    Renvoie `(left, top, width, height, prstGeom_element, (flip_h, flip_v))`
    en EMU natifs (Length de python-pptx, deja consommables tels quels par
    `framed_image.place_image_in_frame`), ou None si aucune shape ne
    correspond. Le flip (`a:xfrm/@flipH|@flipV`) est renvoye a titre
    INFORMATIF seulement — NE PAS le reappliquer aveuglement a une image
    inseree : sur une AUTOSHAPE (fill uni), flipH ne change que le choix du
    coin arrondi (les 2 diagonales sont visuellement identiques pour un
    aplat de couleur) ; sur une IMAGE, flipH retourne aussi le contenu
    pixel — constate au rendu reel sur ce projet (photo mirroir, texte a
    l'envers) apres avoir reproduit le flip source par reflexe. Si le coin
    arrondi doit matcher EXACTEMENT l'original, choisir un `prstGeom`
    equivalent sans toucher au flip de l'image plutot que de flipper
    l'image elle-meme."""
    for sh in shapes:
        spPr = getattr(sh._element, "spPr", None)
        if spPr is None:
            continue
        g = spPr.find(qn("a:prstGeom"))
        if g is None or g.get("prst") != prst:
            continue
        if largeur_min_in is not None and Emu(sh.width).inches < largeur_min_in:
            continue
        xfrm = spPr.find(qn("a:xfrm"))
        flip_h = xfrm is not None and xfrm.get("flipH") == "1"
        flip_v = xfrm is not None and xfrm.get("flipV") == "1"
        return sh.left, sh.top, sh.width, sh.height, g, (flip_h, flip_v)
    return None


def sans_puce(paragraph):
    """Retire l'indentation de puce heritee (marL/indent) et la puce elle-meme
    d'un paragraphe — repris tel quel de VSCode2/VSCode3 (`_sans_puce`).
    Cause reelle du numero de chapitre (« 01 ») qui wrappe sur 2 lignes dans le
    petit encart-pilule du layout Chapitre : le style de liste herite pose
    marL=0.5in dans un encart de ~0.55in de large. python-pptx n'expose pas ces
    attributs -> manipulation XML directe. `buNone` explicite en plus du retrait
    des balises de puce existantes : un caractere de puce residuel peut survivre
    a un niveau que marL/indent seuls ne couvrent pas."""
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("marL", "0")
    pPr.set("indent", "0")
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    pPr.append(pPr.makeelement(qn("a:buNone"), {}))


# ---------------------------------------------------------------------------
# Helpers durcis deck binaire (lecons VSCode4, runs du 2026-07-21)
# ---------------------------------------------------------------------------
# Regle AJOUTER-AVANT-SUPPRIMER : quand on remplace une slide d'un deck binaire,
# creer la nouvelle slide AVANT de supprimer l'ancienne — un delete puis add
# dans le meme cycle reutilise un nom de part (slideN.xml) et produit une
# corruption « Duplicate part name » que python-pptx ne voit pas mais qui rend
# le fichier inouvrable dans PowerPoint (HRESULT 0x80CB4404).


def _normaliser(texte):
    return " ".join(str(texte).split()).casefold()


def trouver_slide_par_titre(prs, titre):
    """Retrouve LA slide dont une shape porte exactement `titre` (comparaison
    normalisee : espaces repliees, casse ignoree) et renvoie (index_0base, slide).

    Pourquoi si strict : les matchers approximatifs ont tous ete pieges sur le
    deck OHC — proximite de position (a repeint la forme voisine), title_of
    (a remonte le kicker au-dessus du titre), recherche dans le corps de texte
    (a matche une slide qui CITAIT le titre cherche). L'egalite stricte sur le
    texte complet d'une shape + l'assertion d'unicite rendent l'erreur bruyante
    au lieu de silencieuse.

    Leve ValueError si zero ou plusieurs slides matchent (dans ce cas, resoudre
    l'ambiguite cote appelant — jamais « prendre la premiere »)."""
    cible = _normaliser(titre)
    matches = []
    for idx, slide in enumerate(prs.slides):
        for sh in slide.shapes:
            if not getattr(sh, "has_text_frame", False):
                continue
            if _normaliser(sh.text_frame.text) == cible:
                matches.append((idx, slide))
                break
    if len(matches) != 1:
        detail = ", ".join(f"slide {i + 1}" for i, _ in matches) or "aucune"
        raise ValueError(
            f"titre {titre!r} : {len(matches)} slide(s) trouvee(s) ({detail}) — "
            "1 exigee (assertion d'unicite)")
    return matches[0]


def supprimer_slide(prs, slide):
    """Suppression SURE d'une slide : retire l'entree de sldIdLst ET lache la
    relation associee (drop_rel). Sans le drop_rel, la part de slide devient
    orpheline : python-pptx (parseur tolerant) ne voit rien, mais PowerPoint
    refuse d'ouvrir le fichier au save suivant (0x80CB4404) — lecon payee 2 fois
    sur le deck OHC. Le slide_id est capture AVANT de toucher la liste (il se
    resout via sldIdLst : le lire apres l'avoir videe leve ValueError)."""
    sid = slide.slide_id
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        if int(sld_id.get("id")) == sid:
            prs.part.drop_rel(sld_id.get(qn("r:id")))
            sld_id_lst.remove(sld_id)
            return
    raise ValueError(f"slide id={sid} absente de sldIdLst")


def clear_slides(prs):
    """Retire TOUTES les slides d'une presentation chargee depuis un template
    (on ne veut heriter que masters/layouts/theme). Meme exigence de drop_rel
    que supprimer_slide : vider sldIdLst sans lacher les relations laisse des
    parts orphelines que PowerPoint refuse ensuite d'ouvrir (constate via
    l'automation COM alors que les tests croyaient le fichier valide)."""
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        prs.part.drop_rel(sld_id.get(qn("r:id")))
        sld_id_lst.remove(sld_id)


def purger_rels_slides_orphelines(prs):
    """Filet anti-corruption avant save() sur un deck binaire retravaille :
    lache toute relation de type slide de la part presentation qui n'est plus
    referencee par sldIdLst (heritee d'une suppression historique faite sans
    drop_rel). Renvoie le nombre de relations purgees (0 = deck sain)."""
    rids_utilises = {s.get(qn("r:id")) for s in prs.slides._sldIdLst}
    purges = 0
    for rid, rel in list(prs.part.rels.items()):
        if rel.reltype == RT.SLIDE and rid not in rids_utilises:
            prs.part.drop_rel(rid)
            purges += 1
    return purges
